# Standard library imports
import os
import time
import traceback
from datetime import datetime, timedelta

# Third-party imports
import pandas as pd
import pytz
import requests
import streamlit as st
from pptx import Presentation

# Local application imports
from src.portioning.portion_controller import MealRecommendation
from src.stickers.shipping_sticker_generator import *
from src.stickers.shipping_sticker_generator_v2 import *
from src.stickers.shipping_sticker_generator_v3 import generate_shipping_stickers_barcode
from src.stickers.dish_sticker_generator_airtable import *
from src.stickers.dish_sticker_generator_barcode import *
from src.generators.one_pager_generator import *
from src.data.store_access import new_database_access
from src.generators.clientservings_excel_output import *
from src.generators.to_make_sheet_generator import *
from src.utils.cancellable import CancellableTask

# Streamlit app
def main():
    db = new_database_access()
    # Current date for filename
    
    st.title('Square Fare Toolkits 🧰')
    # Set the time zone to Eastern Standard Time (EST)
    est = pytz.timezone('US/Eastern')
    current_date_time = datetime.now(est).strftime("%Y%m%d_%H%M")

    #
    st.divider()
    st.header('Refresh Landing Page Cache 🔄')
    st.markdown("⚠️ If you have made any changes to the Ingredients, Dishes, or Variants, you need to refresh the cache to see the changes.")

    # API base URL for the landing page backend
    LANDING_PAGE_API_URL = "https://orders.getsquarefare.com"

    refresh_landing_page_cache_button = st.button("Refresh Landing Page Cache")
    if refresh_landing_page_cache_button:
        try:
            with st.spinner('Refreshing caches...'):
                # Call the refresh endpoint for both caches
                response = requests.post(f"{LANDING_PAGE_API_URL}/api/refresh-cache", timeout=120)
                if response.status_code == 200:
                    st.success("Landing page cache refreshed successfully!")
                    st.info("Please wait for a few minutes for the changes to take effect.")
                    
                else:
                    st.error(f"Cache refresh failed: {response.text}")
        except requests.exceptions.ConnectionError:
            st.error("Could not connect to the landing page server. Is it running?")
        except requests.exceptions.Timeout:
            st.error("Request timed out. The cache refresh may still be in progress.")
        except Exception as e:
            st.error(f"Error refreshing cache: {str(e)}")
    # portion algo trigger
    st.divider()
    st.header('Portion Algorithm 🍽️')
    with st.expander('Expand to see more details'):
        st.markdown("⚠️ Please first confirm :red[Open Orders] are :red[reviewed] and :red[approved]")
        st.markdown("⚠️ Currently :red[only] orders included in :red[this view] will be processed: [Open Orders > Running Portioning](https://airtable.com/appEe646yuQexwHJo/tblxT3Pg9Qh0BVZhM/viwrZHgdsYWnAMhtX?blocks=hide)")

        portion_task = st.session_state.get("portion_task")

        if portion_task is not None and not portion_task.is_done():
            # A run is in progress — show Cancel + live status, poll via rerun.
            progress = st.session_state.get("portion_progress", {})
            elapsed_str = str(timedelta(seconds=int(portion_task.elapsed())))
            cancel_requested = portion_task.is_cancelled()
            done_count = progress.get("done", 0)
            failed_count = progress.get("failed", 0)
            total_count = progress.get("total", 0)
            status = progress.get("status", "Starting…")

            col_cancel, col_status = st.columns([1, 3])
            with col_cancel:
                if cancel_requested:
                    st.button("Cancelling…", disabled=True, key="portion_cancel_disabled")
                else:
                    if st.button("Cancel Portioning", key="portion_cancel"):
                        portion_task.cancel()
                        st.rerun()
            with col_status:
                if total_count:
                    detail = f"{done_count}/{total_count} orders done"
                    if failed_count:
                        detail += f" ({failed_count} failed)"
                    detail += f" — {status}"
                else:
                    detail = status
                if cancel_requested:
                    st.warning(f"Cancelling — waiting for in-flight orders to wrap up… ({elapsed_str}) | {detail}")
                else:
                    st.info(f"Running portioning algorithm… 🕐 {elapsed_str} | {detail}")

            # Poll: sleep briefly then rerun until task finishes.
            time.sleep(2)
            st.rerun()
        else:
            if portion_task is not None and portion_task.is_done():
                elapsed_str = str(timedelta(seconds=int(portion_task.elapsed())))
                if portion_task.error is not None:
                    st.error(f"Portioning failed: {portion_task.error}")
                else:
                    finishedCount, failedCount, failedCases = portion_task.result
                    if portion_task.is_cancelled():
                        st.warning(f"Portioning cancelled after {elapsed_str}. {finishedCount} orders completed before cancellation.")
                    else:
                        st.success(f"{finishedCount} orders completed in {elapsed_str}! ✅")
                    if len(failedCases) > 0:
                        if failedCount > 0:
                            st.error(f"{failedCount} orders failed to process. Please review the following cases:")
                            st.write(failedCases)
                        else:
                            st.error("Portioning stopped half way through, please correct the following cases and re-run the algorithm:")
                            st.write(failedCases)

            if st.button("Yeh! Run Portioning Now"):
                meal_recommendation = MealRecommendation()
                progress = {"status": "Starting…", "done": 0, "failed": 0, "total": 0}
                st.session_state.portion_progress = progress
                task = CancellableTask(meal_recommendation.generate_recommendations_with_thread, progress=progress)
                task.start()
                st.session_state.portion_task = task
                st.rerun()

    # clientservings csv download
    st.header(':green[ClientServings Excel] 📑')
    with st.expander('Expand to see more details'):
        st.markdown("⚠️ Source Table: [Clientservings > For Clientservings Excel Output](https://airtable.com/appEe646yuQexwHJo/tblVwpvUmsTS2Se51/viwgt50kLisz8jx7b?blocks=hide)")
        st.markdown("⚠️ If noticed any issues or missing data, please first check data in source table and then re-run the generator")

        clientservings_task = st.session_state.get("clientservings_task")

        if clientservings_task is not None and not clientservings_task.is_done():
            progress = st.session_state.get("clientservings_progress", {})
            elapsed_str = str(timedelta(seconds=int(clientservings_task.elapsed())))
            done = progress.get("done", 0)
            total = progress.get("total", 0)
            status = progress.get("status", "Starting…")
            detail = f"{done}/{total} dishes — {status}" if total else status
            st.info(f"Generating ClientServings Excel… 🕐 {elapsed_str} | {detail}")
            time.sleep(2)
            st.rerun()
        else:
            if clientservings_task is not None and clientservings_task.is_done():
                elapsed_str = str(timedelta(seconds=int(clientservings_task.elapsed())))
                err = clientservings_task.error
                if err is not None:
                    if isinstance(err, AirTableError):
                        st.error(f"Error generating ClientServings Excel: {err}")
                    else:
                        st.error("An unexpected error occurred. Please check the data in source table and try again.")
                        with st.expander("Technical Error Details (for debugging)"):
                            st.code("".join(traceback.format_exception(type(err), err, err.__traceback__)))
                else:
                    excel_data = clientservings_task.result
                    clientservings_excel_name = f'{current_date_time}_clientservings.xlsx'
                    st.download_button(
                        label="Download ClientServings Excel",
                        data=excel_data,
                        file_name=clientservings_excel_name,
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        key="clientservings_download"
                    )
                    st.success(f"Excel file generated successfully in {elapsed_str}!")

            if st.button("Get ClientServings"):
                progress = {"status": "Starting…", "done": 0, "total": 0}
                st.session_state.clientservings_progress = progress
                task = CancellableTask(lambda cancel_event=None: generate_clientservings_excel(db, progress=progress))
                task.start()
                st.session_state.clientservings_task = task
                st.rerun()


    # shipping_sticker_generator_v1
    st.header(':blue[Shipping Sticker] - Local Uploads 🚚')
    with st.expander('Expand to see more details'):
        st.markdown("⚠️ Source Table: Uploaded Sheet")
        # Display the shipping template file for double-checking
        # File uploader for the new template
        template_path = 'template/Shipping_Sticker_Local_Uploads_Template.pptx'
        shipping_data_path = 'template/shipping_stickers_data_template.csv'

        # Display existing template with download option
        if os.path.exists(template_path):
            with open(template_path, "rb") as template_file:
                st.download_button(
                    label="⬇️ View Existing PowerPoint Template",
                    data=template_file,
                    file_name="Shipping_Sticker_Local_Uploads_Template.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key="download_shipping_local_uploads_template"
                )
        else:
            st.warning(f"⚠️ Default template not found at {template_path}")

        # Upload new template to replace existing one
        # new_shipping_template = st.file_uploader(
        #     "📤 Upload new PowerPoint template to replace existing (optional)",
        #     type="pptx",
        #     key="new_shipping_local_uploads_template",
        #     help="Upload a new PowerPoint template to replace the current one"
        # )

        # if new_shipping_template is not None:
        #     try:
        #         # Validate the uploaded template
        #         test_prs = Presentation(new_shipping_template)
        #         if len(test_prs.slides) == 0:
        #             st.error("❌ The uploaded template has no slides. Please upload a valid template.")
        #         else:
        #             # Save the new template to the template folder
        #             os.makedirs('template', exist_ok=True)
        #             with open(template_path, 'wb') as f:
        #                 f.write(new_shipping_template.getvalue())
        #             st.success(f"✅ Template successfully replaced! New template saved to {template_path}")
        #     except Exception as e:
        #         st.error(f"❌ Error loading uploaded template: {str(e)}")
        #         st.info("Please upload a valid PowerPoint template file")

        # st.divider()

        st.subheader("REQUIRED: upload csv file of address data",divider="blue")
           # Display existing shipping data template with download option
        if os.path.exists(shipping_data_path):
            with open(shipping_data_path, "rb") as template_file:
                st.download_button(
                    label="⬇️ Download Existing Shipping Data Template for Required Format",
                    data=template_file,
                    file_name="shipping_stickers_data_template.csv",
                    mime="text/csv",
                    key="download_shipping_data_template"
                )
        else:
            st.warning(f"⚠️ Default shipping data template not found at {shipping_data_path}")

        uploaded_shipping_file = st.file_uploader(":blue[Upload shipping_stickers_data.csv]", type="csv")
        

        if template_path is not None and uploaded_shipping_file is not None:
            df_shipping = process_shipping_data(uploaded_shipping_file)
            num_slide_per_address = st.number_input("Enter the number of slides per address (Default: 1)", value=1, min_value=1, max_value=100, step=1)
            sticker_generate_button = st.button("Generate Stickers")
            if sticker_generate_button:

                # Generate the PowerPoint file
                prs =generate_ppt_v2(df_shipping, num_slide_per_address,template_path)

                if prs is not None:
                    st.markdown(":green[Stickers are ready! Click to download]")
                    current_date_time = datetime.now(est).strftime("%Y%m%d_%H%M")
                    updated_ppt_name = f'{current_date_time}_shipping_sticker_local_uploads.pptx'
                    prs.save(updated_ppt_name)
                    with open(updated_ppt_name, "rb") as file:
                        st.download_button(
                            label="Download Stickers",
                            data=file,
                            file_name=updated_ppt_name,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )


    # shipping_sticker_generator_v2
    st.header(':blue[Shipping Sticker] - Airtable 🚚')
    with st.expander('Expand to see more details'):
        st.markdown("⚠️ Source Table: [Open Orders > For Shipping Stickers](https://airtable.com/appEe646yuQexwHJo/tblxT3Pg9Qh0BVZhM/viwDpTtU0qaT9NcvG?blocks=hide)")
        st.markdown("⚠️ If noticed any issues or missing data, please first check data in source table and then re-run the generator")

        # Template file management
        template_path = 'template/Shipping_Sticker_Template_v2.pptx'

        # Display existing template with download option
        if os.path.exists(template_path):
            with open(template_path, "rb") as template_file:
                st.download_button(
                    label="⬇️ Download Existing Template",
                    data=template_file,
                    file_name="Shipping_Sticker_Template_v2.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key="download_shipping_template"
                )
        else:
            st.warning(f"⚠️ Default template not found at {template_path}")

        # Upload new template to replace existing one
        # new_shipping_template = st.file_uploader(
        #     "📤 Upload new template to replace existing (optional)",
        #     type="pptx",
        #     key="new_shipping_sticker_template",
        #     help="Upload a new PowerPoint template to replace the current one"
        # )

        # if new_shipping_template is not None:
        #     try:
        #         # Validate the uploaded template
        #         test_prs = Presentation(new_shipping_template)
        #         if len(test_prs.slides) == 0:
        #             st.error("❌ The uploaded template has no slides. Please upload a valid template.")
        #         else:
        #             # Save the new template to the template folder
        #             os.makedirs('template', exist_ok=True)
        #             with open(template_path, 'wb') as f:
        #                 f.write(new_shipping_template.getvalue())
        #             st.success(f"✅ Template successfully replaced! New template saved to {template_path}")
        #     except Exception as e:
        #         st.error(f"❌ Error loading uploaded template: {str(e)}")
        #         st.info("Please upload a valid PowerPoint template file")

        # st.divider()
        shipping_sticker_generate_button = st.button("Generate Shipping Stickers")

        if shipping_sticker_generate_button:
            try:
                with st.spinner('Generating shipping stickers... It may take a few minutes 🕐'):
                    ppt_file = generate_shipping_stickers(db)

                    if ppt_file:
                        updated_ppt_name = f'{current_date_time}_shipping_sticker.pptx'

                        # Provide download button
                        st.download_button(
                            label="Download Stickers",
                            data=ppt_file.getvalue(),
                            file_name=updated_ppt_name,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )

                        st.success(f"✅ Successfully generated shipping stickers!")
                    else:
                        st.warning("⚠️ No stickers were generated. Please check if there are any open orders in Airtable.")
            except AirTableError as e:
                st.error(f"⚠️ Airtable Connection Error: {str(e)}")
                st.info("Please check your Airtable API key and connection settings.")
            except PPTGenerationError as e:
                st.error(f"⚠️ PowerPoint Generation Error: {str(e)}")
                st.info("There was a problem creating the PowerPoint file. Please check your template.")
            except ValueError as e:
                st.warning(f"⚠️ {str(e)}")
            except Exception as e:
                st.error(f"⚠️ Unexpected Error: {str(e)}")
                st.info("Please contact support with the error details and time of occurrence.")
                # Optionally show technical details in an expander
                with st.expander("Technical Error Details (for debugging)"):
                    st.code(traceback.format_exc())

    # shipping_sticker_generator_v3 (Airtable + Barcode)
    st.header(':blue[Shipping Sticker] - Airtable + Barcode 🚚')
    with st.expander('Expand to see more details'):
        st.markdown("⚠️ Source Table: [Open Orders > For Shipping Stickers](https://airtable.com/appEe646yuQexwHJo/tblxT3Pg9Qh0BVZhM/viwDpTtU0qaT9NcvG?blocks=hide)")
        st.markdown("⚠️ If noticed any issues or missing data, please first check data in source table and then re-run the generator")

        template_path_v3 = 'template/Shipping_Sticker_Template_v3.pptx'

        if os.path.exists(template_path_v3):
            with open(template_path_v3, "rb") as template_file:
                st.download_button(
                    label="⬇️ Download Existing Template",
                    data=template_file,
                    file_name="Shipping_Sticker_Template_v3.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key="download_shipping_template_v3"
                )
        else:
            st.warning(f"⚠️ Default template not found at {template_path_v3}")

        shipping_sticker_v3_generate_button = st.button("Generate Shipping Stickers with Barcode")

        if shipping_sticker_v3_generate_button:
            try:
                with st.spinner('Generating shipping stickers with barcode... It may take a few minutes 🕐'):
                    ppt_file, shipping_list = generate_shipping_stickers_barcode(db)

                    if ppt_file:
                        updated_ppt_name = f'{current_date_time}_shipping_sticker_barcode.pptx'

                        st.download_button(
                            label="Download Stickers",
                            data=ppt_file.getvalue(),
                            file_name=updated_ppt_name,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            key="shipping_sticker_v3_download"
                        )

                        st.success(f"✅ Successfully generated shipping stickers with barcode for {len(shipping_list)} unique bag records!")
                    else:
                        st.warning("⚠️ No stickers were generated. Please check if there are any open orders in Airtable.")
            except AirTableError as e:
                st.error(f"⚠️ Airtable Connection Error: {str(e)}")
                st.info("Please check your Airtable API key and connection settings.")
            except PPTGenerationError as e:
                st.error(f"⚠️ PowerPoint Generation Error: {str(e)}")
                st.info("There was a problem creating the PowerPoint file. Please check your template.")
            except ValueError as e:
                st.warning(f"⚠️ {str(e)}")
            except Exception as e:
                st.error(f"⚠️ Unexpected Error: {str(e)}")
                st.info("Please contact support with the error details and time of occurrence.")
                with st.expander("Technical Error Details (for debugging)"):
                    st.code(traceback.format_exc())

    # dish_sticker_generator (Airtable)
    st.header(':orange[Dish Sticker] - Airtable 🍱')
    with st.expander('Expand to see more details'):
        st.markdown("⚠️ Source Table: [Clientservings > Sorted (Dish Sticker)](https://airtable.com/appEe646yuQexwHJo/tblVwpvUmsTS2Se51/viw5hROs9I9vV0YEq?blocks=bipVZAG8G3VXIa12K)")
        st.markdown("⚠️ If noticed any issues or missing data, please first check data in source table and then re-run the generator")

        # Template file handling with error checking
        try:
            # new_dish_sticker_template = st.file_uploader(":blue[Optionally Upload a new Dish_Sticker_Template.pptx]", type="pptx", key="new_dish_sticker_template")
            
            # if new_dish_sticker_template is not None:
            #     try:
            #         prs_file = Presentation(new_dish_sticker_template)
            #         st.success("Custom template loaded successfully!")
            #     except Exception as e:
            #         st.error(f"Error loading uploaded template: {str(e)}")
            #         st.info("Please upload a valid PowerPoint template file")
            #         st.stop()
            # else:
            template_path = 'template/Dish_Sticker_Template.pptx'
            try:
                if not os.path.exists(template_path):
                    st.error(f"Default template not found at {template_path}")
                    st.info("Please upload a custom template to continue")
                    st.stop()
                else:
                    with open(template_path, "rb") as template_file:
                        st.download_button(
                            label="⬇️ View Existing Template",
                            data=template_file,
                            file_name="Dish_Sticker_Template.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            key="download_dish_sticker_template"
                        )
                prs_file = Presentation(template_path)
            except Exception as e:
                st.error(f"Error loading default template: {str(e)}")
                st.info("Please check that the template file exists and is not corrupted")
                st.stop()
        except Exception as e:
            st.error(f"File upload error: {str(e)}")
            st.stop()
        
        # Generate button
        st.markdown("⚠️ Make sure you grouped orders by Dish ID and generate position id in Airtable, before running this generator")
        dish_sticker_generate_button = st.button("I have applied position id, lets get Dish Stickers")
        
        if dish_sticker_generate_button:
            try:
                with st.spinner('Generating dish stickers... It may take a few minutes 🕐'):
                    # Check if the presentation has slides
                    if len(prs_file.slides) == 0:
                        st.error("The template contains no slides. Please use a valid template.")
                        st.stop()
                    
                    # Generate stickers with error handling
                    prs = generate_dish_stickers(prs_file)
                    
                    # Check if generation was successful
                    if prs is None:
                        st.error("Failed to generate stickers. Please check the error messages above.")
                        st.stop()
                    
                    # Save the file with error handling
                    updated_ppt_name = f'{current_date_time}_dish_sticker.pptx'
                    try:
                        prs.save(updated_ppt_name)
                        st.success(f"Successfully saved {updated_ppt_name}")
                    except Exception as e:
                        st.error(f"Failed to save PowerPoint file: {str(e)}")
                        st.info("Check if you have write permissions to the current directory")
                        st.stop()
                
                # Provide download button with error handling
                try:
                    with open(updated_ppt_name, "rb") as file:
                        st.download_button(
                            label="Download Stickers",
                            data=file,
                            file_name=updated_ppt_name,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                except Exception as e:
                    st.error(f"Error preparing download: {str(e)}")
                    st.info(f"The file was saved as {updated_ppt_name} but cannot be downloaded directly. Please access it from your server.")
            except Exception as e:
                st.error(f"Unexpected error during sticker generation: {str(e)}")
                st.exception(e)  

    # Barcode Dish Sticker Generator
    st.header(':orange[Dish Sticker] - Barcode 🍱')
    with st.expander('Expand to see more details'):
        st.markdown("⚠️ Source Table: [Clientservings > Sorted (Dish Sticker)](https://airtable.com/appEe646yuQexwHJo/tblVwpvUmsTS2Se51/viw5hROs9I9vV0YEq?blocks=bipVZAG8G3VXIa12K)")
        st.markdown("⚠️ If noticed any issues or missing data, please first check data in source table and then re-run the generator")
        
        # Template file handling with error checking for Barcode generator
        try:
            # new_qr_dish_sticker_template = st.file_uploader(":blue[Optionally Upload a new Dish_Sticker_Template_Barcode.pptx]", type="pptx", key="new_qr_dish_sticker_template")
            
            # if new_qr_dish_sticker_template is not None:
            #     try:
            #         qr_prs_file = Presentation(new_qr_dish_sticker_template)
            #         st.success("Custom Barcode template loaded successfully!")
            #     except Exception as e:
            #         st.error(f"Error loading uploaded Barcode template: {str(e)}")
            #         st.info("Please upload a valid PowerPoint template file")
            #         st.stop()
            # else:
            qr_template_path = 'template/Dish_Sticker_Template_Barcode.pptx'
            try:
                if not os.path.exists(qr_template_path):
                    st.error(f"Default Barcode template not found at {qr_template_path}")
                    st.info("Please upload a custom template to continue")
                    st.stop()
                else:
                    with open(qr_template_path, "rb") as template_file:
                        st.download_button(
                            label="⬇️ View Existing Template",
                            data=template_file,
                            file_name="Dish_Sticker_Template_Barcode.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            key="download_dish_sticker_barcode_template"
                        )
                qr_prs_file = Presentation(qr_template_path)
            except Exception as e:
                st.error(f"Error loading default Barcode template: {str(e)}")
                st.info("Please check that the template file exists and is not corrupted")
                st.stop()
        except Exception as e:
            st.error(f"File upload error: {str(e)}")
            st.stop()

        # Generate button for Barcode stickers
        st.markdown("⚠️ Make sure you grouped orders by Dish ID and generate position id in Airtable, before running this generator")

        barcode_task = st.session_state.get("barcode_task")

        if barcode_task is not None and not barcode_task.is_done():
            # Running — show Cancel + live status, poll via rerun.
            progress = st.session_state.get("barcode_progress", {})
            elapsed_str = str(timedelta(seconds=int(barcode_task.elapsed())))
            cancel_requested = barcode_task.is_cancelled()
            slide_count = progress.get("slide_count", 0)
            total_slides = progress.get("total_slides", 0)
            status = progress.get("status", "Starting…")

            col_cancel, col_status = st.columns([1, 3])
            with col_cancel:
                if cancel_requested:
                    st.button("Cancelling…", disabled=True, key="barcode_cancel_disabled")
                else:
                    if st.button("Cancel Barcode Stickers", key="barcode_cancel"):
                        barcode_task.cancel()
                        st.rerun()
            with col_status:
                detail = f"{slide_count}/{total_slides} slides — {status}" if total_slides else status
                if cancel_requested:
                    st.warning(f"Cancelling — finishing current slide… ({elapsed_str}) | {detail}")
                else:
                    st.info(f"Generating Barcode dish stickers… 🕐 {elapsed_str} | {detail}")

            time.sleep(2)
            st.rerun()
        else:
            if barcode_task is not None and barcode_task.is_done():
                elapsed_str = str(timedelta(seconds=int(barcode_task.elapsed())))
                if barcode_task.error is not None:
                    st.error(f"Unexpected error during Barcode sticker generation: {barcode_task.error}")
                    with st.expander("Technical Error Details (for debugging)"):
                        st.code("".join(traceback.format_exception(type(barcode_task.error), barcode_task.error, barcode_task.error.__traceback__)))
                else:
                    prs = barcode_task.result
                    # Save the presentation once per completed task.
                    saved_path = st.session_state.get("barcode_saved_path")
                    if saved_path is None and prs is not None:
                        saved_path = f'{current_date_time}_dish_sticker_barcode.pptx'
                        prs.save(saved_path)
                        st.session_state.barcode_saved_path = saved_path

                    progress = st.session_state.get("barcode_progress", {})
                    slide_count = progress.get("slide_count", 0)
                    total_slides = progress.get("total_slides", 0)

                    if barcode_task.is_cancelled():
                        st.warning(f"Cancelled after {elapsed_str} — {slide_count}/{total_slides} slides completed. Partial file saved.")
                    else:
                        st.success(f"Successfully generated Barcode dish stickers in {elapsed_str}!")

                    if saved_path is not None and os.path.exists(saved_path):
                        with open(saved_path, "rb") as file:
                            st.download_button(
                                label="Download Barcode Stickers",
                                data=file,
                                file_name=saved_path,
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                key="barcode_download"
                            )

            if st.button("I have applied position id, lets get Barcode Dish Stickers"):
                if len(qr_prs_file.slides) == 0:
                    st.error("The Barcode template contains no slides. Please use a valid template.")
                    st.stop()
                # Clear prior task state so a fresh run doesn't reuse old saved file / progress.
                st.session_state.pop("barcode_saved_path", None)
                progress = {"status": "Starting…", "slide_count": 0, "total_slides": 0}
                st.session_state.barcode_progress = progress
                task = CancellableTask(generate_dish_stickers_barcode, db, progress=progress)
                task.start()
                st.session_state.barcode_task = task
                st.rerun()

    # one_pager_generator
    st.header('One-Sheeter - Airtable 📑')
    with st.expander('Expand to see more details'):
        st.markdown("⚠️ Source Table: [Open Orders > For One-Sheeter](https://airtable.com/appEe646yuQexwHJo/tblxT3Pg9Qh0BVZhM/viwuVy9aN2LLZrcPF?blocks=hide)")
        st.markdown("⚠️ If noticed any issues or missing data, please first check data in source table and then re-run the generator")

        # File upload for template
        # new_one_pager_template = st.file_uploader(":blue[Optionally Upload new One_Pager_Template.pptx]", 
        #                                         type="pptx", 
        #                                         key="new_one_pager_template",
        #                                         help="Include instruction slide as second slide")

        st.markdown("⚠️ Please make sure all meal stickers are generated before running this one-sheeter generator")

        one_sheeter_task = st.session_state.get("one_sheeter_task")

        if one_sheeter_task is not None and not one_sheeter_task.is_done():
            progress = st.session_state.get("one_sheeter_progress", {})
            elapsed_str = str(timedelta(seconds=int(one_sheeter_task.elapsed())))
            done = progress.get("done", 0)
            total = progress.get("total", 0)
            status = progress.get("status", "Starting…")
            detail = f"{done}/{total} pages — {status}" if total else status
            st.info(f"Generating One-Sheeter… 🕐 {elapsed_str} | {detail}")
            time.sleep(2)
            st.rerun()
        else:
            if one_sheeter_task is not None and one_sheeter_task.is_done():
                elapsed_str = str(timedelta(seconds=int(one_sheeter_task.elapsed())))
                err = one_sheeter_task.error
                if err is not None:
                    st.error(f"Error generating One-Sheeter: {err}")
                    with st.expander("Technical Error Details (for debugging)"):
                        st.code("".join(traceback.format_exception(type(err), err, err.__traceback__)))
                else:
                    prs = one_sheeter_task.result
                    saved_path = st.session_state.get("one_sheeter_saved_path")
                    if saved_path is None and prs is not None:
                        saved_path = f'{current_date_time}_one_sheeter.pptx'
                        prs.save(saved_path)
                        st.session_state.one_sheeter_saved_path = saved_path

                    if saved_path is not None and os.path.exists(saved_path):
                        with open(saved_path, "rb") as file:
                            st.download_button(
                                label="Download One-Sheeter",
                                data=file,
                                file_name=saved_path,
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                key="one_sheeter_download"
                            )
                    st.success(f"One-Sheeter generated successfully in {elapsed_str}!")

            if st.button("Yes I ran the meal sticker already, now lets generate One-Sheeter"):
                template_path = 'template/One_Pager_Template_v2.pptx'
                prs_file = Presentation(template_path)
                if len(prs_file.slides) < 2:
                    st.warning("Your template should include an instruction slide as the second slide. Proceeding without instructions.")
                st.session_state.pop("one_sheeter_saved_path", None)
                progress = {"status": "Starting…", "done": 0, "total": 0}
                st.session_state.one_sheeter_progress = progress
                task = CancellableTask(lambda cancel_event=None: generate_one_pagers(db, template_path, progress=progress))
                task.start()
                st.session_state.one_sheeter_task = task
                st.rerun()

    # to_make_sheet_generator
    st.header('To-Make Sheet Generator 🍳')
    with st.expander('Expand to see more details'):
        st.markdown("⚠️ Source Table: [Clientservings > For To-Make Sheet](https://airtable.com/appEe646yuQexwHJo/tblVwpvUmsTS2Se51/viw4WN1XsjMnHwMkt?blocks=hide)")

        to_make_sheet_task = st.session_state.get("to_make_sheet_task")

        if to_make_sheet_task is not None and not to_make_sheet_task.is_done():
            elapsed_str = str(timedelta(seconds=int(to_make_sheet_task.elapsed())))
            st.info(f"Generating To-Make Sheet… 🕐 {elapsed_str}")
            time.sleep(2)
            st.rerun()
        else:
            if to_make_sheet_task is not None and to_make_sheet_task.is_done():
                elapsed_str = str(timedelta(seconds=int(to_make_sheet_task.elapsed())))
                err = to_make_sheet_task.error
                if err is not None:
                    st.error(f"Error generating to-make sheet: {err}")
                    st.info("Please check the source table data and try again")
                    with st.expander("Technical Error Details (for debugging)"):
                        st.code("".join(traceback.format_exception(type(err), err, err.__traceback__)))
                else:
                    excel_file = to_make_sheet_task.result
                    updated_xlsx_name = f'{current_date_time}_to_make_sheet.xlsx'
                    st.download_button(
                        label="Download To-Make Sheet",
                        data=excel_file.getvalue(),
                        file_name=updated_xlsx_name,
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        key="to_make_sheet_download"
                    )
                    st.success(f"To-Make Sheet generated successfully in {elapsed_str}!")

            if st.button("Generate To-Make Sheet"):
                task = CancellableTask(lambda cancel_event=None: generate_to_make_sheet(db))
                task.start()
                st.session_state.to_make_sheet_task = task
                st.rerun()

if __name__ == "__main__":
    main()
