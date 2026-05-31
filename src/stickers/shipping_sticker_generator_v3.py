import sys
import copy
import hashlib
import logging
from io import BytesIO
from pathlib import Path
from datetime import datetime

import barcode
from barcode.writer import ImageWriter
from PIL import Image, ImageChops
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_FILL
BASE_DIR = Path(__file__).resolve().parents[2]  # Streamlit-App
sys.path.append(str(BASE_DIR))

from src.data.exceptions import AirTableError
from src.data.store_access import new_database_access
from src.stickers.dish_barcode_ids import dish_barcode_from_open_order_fields


logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger(__name__)

VIEW = "viwDpTtU0qaT9NcvG"

DEFAULT_BAG_TEMPLATE = BASE_DIR / "template" / "Shipping_Sticker_Template_v3.pptx"
SNOWFLAKE_IMAGE = BASE_DIR / "template" / "snowflake.png"
SNOWFLAKE_GLYPH = "❄"
NAME_FONT_PT = 40
ADDRESS_FONT_PT = 26
ZONE_FONT_PT = 30
ADDRESS_ROW_HEIGHT = Inches(0.62)
CITY_ROW_HEIGHT = Inches(0.42)
PHONE_ROW_HEIGHT = Inches(0.42)
ADDRESS_STACK_ROW_GAP = Inches(0.02)
NAME_ROW_NUDGE_DOWN = Inches(0.14)
CONTENT_NUDGE_DOWN = Inches(0.10)
NAME_ADDRESS_GAP = Inches(0.18)


class PPTGenerationError(Exception):
    pass


def unwrap(value, default=""):
    if isinstance(value, list):
        return value[0] if value else default
    return value if value is not None else default


def format_phone(phone):
    if not phone:
        return ""

    phone_digits = "".join(filter(str.isdigit, str(phone)))
    if len(phone_digits) == 10:
        return f"{phone_digits[:3]}-{phone_digits[3:6]}-{phone_digits[6:]}"
    return str(phone)


def make_bag_group_key(record):
    return (
        str(record.get("Delivery Date", "")).strip(),
        str(record.get("Shipping Name", "")).strip().lower(),
        str(record.get("Shipping Address 1", "")).strip().lower(),
        str(record.get("Shipping Address 2", "")).strip().lower(),
        str(record.get("Shipping City", "")).strip().lower(),
        str(record.get("Shipping Province", "")).strip().lower(),
        str(record.get("Shipping Postal Code", "")).strip(),
        str(record.get("Zone Number", "")).strip(),
    )


def customization_tags_from_fields(fields):
    """
    Airtable field names vary (e.g. 'Customization Tags' vs
    'Customization Tags (from To_Match_Client_Nutrition)').
    Merge string values from every field whose name contains 'Customization Tags'.
    """
    parts = []
    for key, val in (fields or {}).items():
        if "Customization Tags" not in str(key):
            continue
        if isinstance(val, list):
            parts.extend(str(x) for x in val if str(x).strip())
        elif val is not None and str(val).strip():
            parts.append(str(val))
    return " ".join(parts)


def make_bag_barcode(shipping_info, chunk_index=0):
    delivery_date = str(shipping_info.get("Delivery Date", "")).replace("-", "")
    if not delivery_date:
        delivery_date = datetime.now().strftime("%Y%m%d")

    zone = str(shipping_info.get("Zone Number", "X")).replace("ZONE", "").replace("Zone", "").strip()

    name_addr = "|".join([
        str(shipping_info.get("Shipping Name", "")).strip().lower(),
        str(shipping_info.get("Shipping Address 1", "")).strip().lower(),
        str(shipping_info.get("Shipping Address 2", "")).strip().lower(),
        str(shipping_info.get("Shipping Postal Code", "")).strip(),
        str(chunk_index),
    ])
    short_hash = hashlib.sha1(name_addr.encode("utf-8")).hexdigest()[:6].upper()

    return f"BAG-{delivery_date}-Z{zone}-{short_hash}"


def process_order_data(db):
    """
    Read open orders and group them into shipping/bag records.
    """
    logger.info("Starting shipping sticker barcode data processing")

    orders = db.get_all_open_orders(view=VIEW)

    if not orders:
        raise AirTableError("No open orders found")

    logger.info(f"Found {len(orders)} open orders")

    shipping_data = []

    for order in orders:
        fields = order.get("fields", {})

        required_fields = ["Shipping Name", "Shipping Address 1", "Quantity"]
        if not all(fields.get(field) for field in required_fields):
            logger.warning(f"Skipping order with missing fields: {order.get('id', 'unknown')}")
            continue

        meal_portion = unwrap(fields.get("Meal Portion", ""))
        quantity = unwrap(fields.get("Quantity", 0), 0)

        try:
            quantity = float(quantity)
        except Exception:
            quantity = 0

        if meal_portion == "Breakfast":
            adjusted_quantity = quantity * 0.8
        elif meal_portion == "Snack":
            adjusted_quantity = quantity * 0.1
        else:
            adjusted_quantity = quantity

        zone = unwrap(fields.get("Zone Number (from Delivery Zone)", "N/A"))
        delivery_date = unwrap(fields.get("Delivery Date", ""))

        dish_barcode = dish_barcode_from_open_order_fields(fields)

        shipping_record = {
            "Delivery Date": delivery_date,
            "Shipping Name": unwrap(fields.get("Shipping Name", "")),
            "Shipping Address 1": unwrap(fields.get("Shipping Address 1", "")),
            "Shipping Address 2": unwrap(fields.get("Shipping Address 2", "")),
            "Shipping City": unwrap(fields.get("Shipping City", "")),
            "Shipping Province": str(unwrap(fields.get("Shipping Province", ""))).upper(),
            "Shipping Postal Code": unwrap(fields.get("Shipping Postal Code", "")),
            "Shipping Phone": format_phone(unwrap(fields.get("Shipping Phone", ""))),
            "Zone Number": str(zone),
            "Quantity": adjusted_quantity,

            # Mapping data for bag scan: same id as Open Orders "Portion Result (in ClientServings)" (Client Servings record id).
            "Dish Barcode": dish_barcode,
            "Customer Name": unwrap(fields.get("Customer Name", "")),
            "Meal Portion": meal_portion,
            "Meal Sticker": unwrap(fields.get("Meal Sticker", "")),
            "Customization Tags": customization_tags_from_fields(fields),
        }

        shipping_data.append(shipping_record)

    grouped_shipping = {}

    for record in shipping_data:
        key = make_bag_group_key(record)

        if key not in grouped_shipping:
            grouped_shipping[key] = {
                "Delivery Date": record["Delivery Date"],
                "Shipping Name": record["Shipping Name"],
                "Shipping Address 1": record["Shipping Address 1"],
                "Shipping Address 2": record["Shipping Address 2"],
                "Shipping City": record["Shipping City"],
                "Shipping Province": record["Shipping Province"],
                "Shipping Postal Code": record["Shipping Postal Code"],
                "Shipping Phone": record["Shipping Phone"],
                "Zone Number": record["Zone Number"],
                "Total Quantity": 0,
                "Household Members": set(),
                "Dishes": [],
                "Ice Pack Required": False,
            }

        grouped_shipping[key]["Total Quantity"] += record["Quantity"]

        if record["Customer Name"]:
            grouped_shipping[key]["Household Members"].add(record["Customer Name"])

        dish_line = " - ".join(
            part for part in [
                record.get("Customer Name", ""),
                record.get("Meal Portion", ""),
                record.get("Meal Sticker", ""),
            ]
            if str(part).strip()
        )

        if dish_line:
            grouped_shipping[key]["Dishes"].append(
                {
                    "dishBarcode": record.get("Dish Barcode", ""),
                    "customerName": record.get("Customer Name", ""),
                    "mealPortion": record.get("Meal Portion", ""),
                    "mealSticker": record.get("Meal Sticker", ""),
                    "displayText": dish_line,
                    "status": "unscanned",
                    "adjustedQuantity": record["Quantity"],
                }
            )

        if "Ice Pack" in str(record.get("Customization Tags", "")):
            grouped_shipping[key]["Ice Pack Required"] = True

    shipping_list = []
    portion_per_bag = 6.8

    for shipping_info in grouped_shipping.values():
        dishes = shipping_info["Dishes"]
        household_members = sorted(list(shipping_info["Household Members"]))

        chunks = []
        current_chunk = []
        current_qty = 0.0
        for dish in dishes:
            dish_qty = float(dish.get("adjustedQuantity", 0) or 0)
            if current_chunk and current_qty + dish_qty > portion_per_bag:
                chunks.append(current_chunk)
                current_chunk = []
                current_qty = 0.0
            current_chunk.append(dish)
            current_qty += dish_qty
        if current_chunk:
            chunks.append(current_chunk)
        if not chunks:
            chunks = [[]]

        for chunk_index, chunk_dishes in enumerate(chunks):
            bag = {
                **shipping_info,
                "Dishes": chunk_dishes,
                "Household Members": household_members,
                "Stickers Needed": 1,
                "Bag Barcode": make_bag_barcode(shipping_info, chunk_index),
            }
            shipping_list.append(bag)

    logger.info(f"Processed {len(shipping_list)} bag records")

    return shipping_list


def _style_sticker_paragraph(paragraph, size_pt):
    paragraph.font.size = Pt(size_pt)
    paragraph.font.name = "Lato"


def _shape_text(shape):
    if not shape.has_text_frame:
        return ""
    return shape.text_frame.text or ""


_STICKER_SHAPE_TAGS = {
    "name": "stickerShippingName",
    "address": "stickerShippingAddress",
    "city": "stickerShippingCity",
    "phone": "stickerShippingPhone",
    "zone": "stickerZone",
}


def _sticker_text_shapes(slide):
    found = {}
    tagged = {v: k for k, v in _STICKER_SHAPE_TAGS.items()}
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        if sh.name in tagged:
            found[tagged[sh.name]] = sh
    if found:
        return found

    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        text = _shape_text(sh)
        name = sh.name

        if name == "shippingName" or "Shipping Name" in text:
            found["name"] = sh
        elif name == "shippingAddress" or (
            "Address" in text and "City" not in text and "Phone" not in text
        ):
            found["address"] = sh
        elif name == "shippingZip" or ("City" in text and "Province" in text):
            found["city"] = sh
        elif name == "shippingPhone" or "Shipping Phone" in text:
            found["phone"] = sh
        elif name == "zone" or text.strip().upper() == "ZONE":
            found["zone"] = sh

    return found


def _tag_sticker_layout_shapes(slide):
    """Pin stable names on template shapes before placeholder text is replaced."""
    layout = _sticker_text_shapes(slide)
    template_geom = {
        key: (int(sh.left), int(sh.top), int(sh.width), int(sh.height))
        for key, sh in layout.items()
    }
    for key, sh in layout.items():
        tag = _STICKER_SHAPE_TAGS.get(key)
        if tag and sh is not None:
            sh.name = tag
    return layout, template_geom


def _barcode_horizontal_bounds(slide):
    """Address left edge through ZONE column right edge."""
    layout = _sticker_text_shapes(slide)
    addr = layout.get("address")
    _, _, zone_right = _zone_column_bounds(slide)
    left = int(addr.left) if addr is not None else int(Inches(0.791))
    return left, int(zone_right) - left


def _barcode_max_bottom(slide):
    """Top border of REFRIGERATE/ZONE box — barcode must stay above this line."""
    return _footer_box_line_top(slide) - int(BARCODE_GAP_ABOVE_FOOTER)


def _text_column_bounds(slide):
    layout = _sticker_text_shapes(slide)
    addr = layout.get("address")
    name = layout.get("name")
    if addr is not None:
        return int(addr.left), int(addr.width)
    left = int(name.left if name is not None else Inches(0.789))
    _, _, zone_right = _zone_column_bounds(slide)
    return left, int(zone_right) - left


def _send_shape_to_back(slide, shape):
    el = shape.element
    tree = slide.shapes._spTree
    tree.remove(el)
    tree.insert(2, el)


def _transparent_text_box(shape):
    try:
        shape.fill.background()
    except Exception:
        pass
    try:
        shape.line.fill.background()
    except Exception:
        pass


def _apply_text_frame_left(shape, vertical_anchor=MSO_ANCHOR.MIDDLE):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = vertical_anchor
    for paragraph in tf.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT


def _find_all_white_blocks(slide):
    blocks = []
    tagged = set(_STICKER_SHAPE_TAGS.values())
    for sh in slide.shapes:
        if sh.name in tagged:
            continue
        if not sh.has_text_frame or _shape_text(sh).strip():
            continue
        try:
            if sh.fill.type == MSO_FILL.SOLID:
                blocks.append(sh)
        except Exception:
            continue
    return blocks


def _send_white_blocks_to_back(slide):
    for block in _find_all_white_blocks(slide):
        _send_shape_to_back(slide, block)


def _bring_shape_to_front(slide, shape):
    el = shape.element
    tree = slide.shapes._spTree
    tree.remove(el)
    tree.insert_element_before(el, "p:extLst")


def _find_name_white_block(slide):
    blocks = []
    tagged = set(_STICKER_SHAPE_TAGS.values())
    for sh in slide.shapes:
        if not sh.has_text_frame or sh.name in tagged:
            continue
        if _shape_text(sh).strip():
            continue
        try:
            if sh.fill.type != MSO_FILL.SOLID:
                continue
        except Exception:
            continue
        blocks.append(sh)
    if not blocks:
        return None
    return min(blocks, key=lambda s: (s.top, -s.width))


def _sticker_content_bounds(slide):
    shapes = _sticker_text_shapes(slide)
    left = None
    right = 0
    for key in ("name", "address", "city", "phone"):
        sh = shapes.get(key)
        if sh is None:
            continue
        left = sh.left if left is None else min(left, sh.left)
        right = max(right, sh.left + sh.width)
    if left is None:
        left = Inches(0.79)
    if right <= left:
        right = left + Inches(9.02)
    return int(left), int(right)


def _find_name_shape(slide):
    return _sticker_text_shapes(slide).get("name")


def _zone_column_bounds(slide):
    zone = _sticker_text_shapes(slide).get("zone")
    if zone is not None:
        left = int(zone.left)
        width = int(zone.width)
        return left, width, left + width
    left = int(Inches(7.05))
    width = int(Inches(1.97))
    return left, width, left + width


def _reserve_name_line_for_snowflake(name_shape, zone_left):
    gap = int(Inches(0.1))
    max_width = int(zone_left) - int(name_shape.left) - gap
    if max_width > 0 and name_shape.width > max_width:
        name_shape.width = max_width


def _add_ice_snowflake(slide, name_shape):
    zone_left, zone_width, zone_right = _zone_column_bounds(slide)
    row_top = int(name_shape.top)
    row_height = int(name_shape.height)
    size = int(min(zone_width, row_height, Inches(0.95)))
    pic_left = int(zone_right - size)
    pic_top = row_top

    if SNOWFLAKE_IMAGE.is_file():
        pic = slide.shapes.add_picture(str(SNOWFLAKE_IMAGE), pic_left, pic_top, width=size, height=size)
        pic.name = "iceSnowflake"
        return pic

    box = slide.shapes.add_textbox(zone_left, row_top, zone_width, row_height)
    box.name = "iceSnowflake"
    tf = box.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = SNOWFLAKE_GLYPH
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(96)
    p.font.bold = True
    return box


def _layout_for_copy(template_slide, target_prs):
    """
    Layout must come from target_prs. Using template_slide.slide_layout from another
    Presentation corrupts the saved file (duplicate OPC parts / broken relationships).
    """
    name = template_slide.slide_layout.name
    for layout in target_prs.slide_layouts:
        if layout.name == name:
            return layout
    return target_prs.slide_layouts[0]


def copy_slide(template_slide, target_prs):
    new_slide = target_prs.slides.add_slide(_layout_for_copy(template_slide, target_prs))

    for shape in new_slide.shapes:
        sp = shape._element
        sp.getparent().remove(sp)

    for shape in template_slide.shapes:
        # Copy XML element only — deepcopy(shape) breaks on python-pptx Fill objects (3.13+).
        new_element = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_element, "p:extLst")

    return new_slide


def populate_sticker(slide, shipping_info, ice_pack_required=False):
    layout, template_geom = _tag_sticker_layout_shapes(slide)

    name_shape = layout.get("name")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        text_frame = shape.text_frame

        for paragraph in text_frame.paragraphs:
            text = paragraph.text

            if shape is layout.get("name") or "Shipping Name" in text:
                paragraph.text = shipping_info["Shipping Name"]
                _style_sticker_paragraph(paragraph, NAME_FONT_PT)
                name_shape = shape

            elif shape is layout.get("address") or (
                "Address" in text and "City" not in text and "Phone" not in text
            ):
                if shipping_info["Shipping Address 2"]:
                    paragraph.text = (
                        f"{shipping_info['Shipping Address 1']}, "
                        f"{shipping_info['Shipping Address 2']}"
                    )
                else:
                    paragraph.text = shipping_info["Shipping Address 1"]
                _style_sticker_paragraph(paragraph, ADDRESS_FONT_PT)

            elif shape is layout.get("city") or ("City" in text and "Province" in text):
                paragraph.text = (
                    f"{shipping_info['Shipping City']}, "
                    f"{shipping_info['Shipping Province']} "
                    f"{shipping_info['Shipping Postal Code']}"
                )
                _style_sticker_paragraph(paragraph, ADDRESS_FONT_PT)

            elif shape is layout.get("phone") or "Shipping Phone" in text:
                paragraph.text = str(shipping_info["Shipping Phone"])
                _style_sticker_paragraph(paragraph, ADDRESS_FONT_PT)

            elif shape is layout.get("zone") or text.strip().upper() == "ZONE":
                paragraph.text = f"ZONE {shipping_info['Zone Number']}"
                _style_sticker_paragraph(paragraph, ZONE_FONT_PT)

    layout = _sticker_text_shapes(slide)
    _layout_sticker_text_stack(slide, layout, template_geom)
    name_shape = layout.get("name") or name_shape

    if ice_pack_required and name_shape:
        zone_left, _, _ = _zone_column_bounds(slide)
        _reserve_name_line_for_snowflake(name_shape, zone_left)
        snowflake = _add_ice_snowflake(slide, name_shape)
        _bring_shape_to_front(slide, name_shape)
        if snowflake is not None:
            _bring_shape_to_front(slide, snowflake)


ADDRESS_STACK_ORDER = ("shippingName", "shippingAddress", "shippingZip", "shippingPhone")
FOOTER_ZONE_TOP = Inches(6.06)
FOOTER_BOX_LINE_OFFSET_ABOVE_ZONE = Inches(0.72)
BARCODE_GAP_ABOVE_FOOTER = Inches(0.04)
BARCODE_GAP_BELOW_ADDRESS = Inches(0.06)
BARCODE_NUDGE_DOWN = Inches(0.10)
BARCODE_IMAGE_HEIGHT = Inches(1.55)
BARCODE_RENDER_DPI = 300
BARCODE_TEXT_HEIGHT = Inches(0.22)
BARCODE_LABEL_FONT_PT = 12


def _emu_to_mm(value_emu):
    """Convert PowerPoint EMU to millimeters (python-barcode uses mm internally)."""
    return value_emu * 25.4 / 914400


def _code128_module_count(barcode_value):
    """Total module units in the encoded bar pattern (for width-fit rendering)."""
    CODE128 = barcode.get_barcode_class("code128")
    writer = ImageWriter()
    obj = CODE128(barcode_value, writer=writer)
    line = obj.build()[0]
    return len(line)


def _trim_image_whitespace(img):
    """Remove built-in quiet-zone padding so bars reach the image edges."""
    bg = Image.new("RGBA", img.size, (255, 255, 255, 255))
    bbox = ImageChops.difference(img, bg).getbbox()
    return img.crop(bbox) if bbox else img


def _render_barcode_bars_png(barcode_value, width_emu, height_emu):
    """
    Render Code128 bars only; module width is computed so bars fill the slot edge-to-edge.
    Human-readable text is added separately as a PPT text box for true alignment.
    """
    target_width_mm = _emu_to_mm(width_emu)
    target_height_mm = _emu_to_mm(height_emu)
    modules = _code128_module_count(barcode_value)
    module_width = target_width_mm / modules

    CODE128 = barcode.get_barcode_class("code128")
    obj = CODE128(barcode_value, writer=ImageWriter(mode="RGBA", dpi=BARCODE_RENDER_DPI))
    buf = BytesIO()
    obj.write(
        buf,
        {
            "quiet_zone": 0,
            "write_text": False,
            "module_height": target_height_mm,
            "module_width": module_width,
            "margin_top": 0,
            "margin_bottom": 0,
            "background": "white",
            "foreground": "black",
        },
    )
    buf.seek(0)
    img = _trim_image_whitespace(Image.open(buf).convert("RGBA"))

    target_w_px = max(1, round(width_emu * BARCODE_RENDER_DPI / 914400))
    target_h_px = max(1, round(height_emu * BARCODE_RENDER_DPI / 914400))
    if img.size != (target_w_px, target_h_px):
        img = img.resize((target_w_px, target_h_px), Image.Resampling.LANCZOS)

    out = BytesIO()
    img.save(out, format="PNG")
    out.seek(0)
    return out


def _constrain_and_reflow_address(slide):
    layout = _sticker_text_shapes(slide)
    template_geom = {
        key: (int(sh.left), int(sh.top), int(sh.width), int(sh.height))
        for key, sh in layout.items()
    }
    _layout_sticker_text_stack(slide, layout, template_geom)


def _row_height_for_key(key):
    if key == "address":
        return int(ADDRESS_ROW_HEIGHT)
    if key == "city":
        return int(CITY_ROW_HEIGHT)
    if key == "phone":
        return int(PHONE_ROW_HEIGHT)
    return None


def _layout_sticker_text_stack(slide, layout, template_geom):
    """
    Name on white strip; address/city/phone stacked compactly below name.
    Never pull rows up into the name row — compact stack leaves room for a tall barcode.
    """
    white = _find_name_white_block(slide)
    name = layout.get("name")
    nudge = int(CONTENT_NUDGE_DOWN)
    cursor = 0

    if white and name:
        nl, _, nw, _ = template_geom.get("name", (int(Inches(0.789)), 0, int(Inches(7.12)), 0))
        name.top = int(white.top) + int(NAME_ROW_NUDGE_DOWN) + nudge
        name.height = int(white.height)
        name.left = nl
        name.width = nw
        _transparent_text_box(name)
        _apply_text_frame_left(name, vertical_anchor=MSO_ANCHOR.MIDDLE)
        cursor = int(name.top + name.height) + int(NAME_ADDRESS_GAP)

    for key in ("address", "city", "phone"):
        sh = layout.get(key)
        if sh is None or key not in template_geom:
            continue
        left, _, width, _ = template_geom[key]
        row_h = _row_height_for_key(key)
        sh.left = left
        sh.top = cursor
        sh.width = width
        if row_h:
            sh.height = row_h
        _transparent_text_box(sh)
        _apply_text_frame_left(sh, vertical_anchor=MSO_ANCHOR.TOP)
        cursor = int(sh.top + sh.height) + int(ADDRESS_STACK_ROW_GAP)

    _send_white_blocks_to_back(slide)
    if name:
        _bring_shape_to_front(slide, name)


def _compute_barcode_rect(slide):
    """Large barcode: address-left to zone-right, bottom anchored above footer box."""
    content_left, content_width = _barcode_horizontal_bounds(slide)
    barcode_max_bottom = _barcode_max_bottom(slide)
    height = int(BARCODE_IMAGE_HEIGHT)
    top = barcode_max_bottom - height
    phone_bottom = _shipping_phone_bottom(slide)
    min_top = phone_bottom + int(BARCODE_GAP_BELOW_ADDRESS)
    if top < min_top:
        top = min_top
    top += int(BARCODE_NUDGE_DOWN)
    if top + height > barcode_max_bottom:
        height = barcode_max_bottom - top
    return content_left, content_width, top, height


def _address_stack_bottom(slide) -> int:
    layout = _sticker_text_shapes(slide)
    bottom = 0
    for key in ("name", "address", "city", "phone"):
        sh = layout.get(key)
        if sh is not None:
            bottom = max(bottom, int(sh.top + sh.height))
    return bottom


def _footer_zone_top(slide):
    zone = _sticker_text_shapes(slide).get("zone")
    if zone is not None:
        return zone.top
    return FOOTER_ZONE_TOP


def _footer_box_line_top(slide):
    return int(_footer_zone_top(slide) - FOOTER_BOX_LINE_OFFSET_ABOVE_ZONE)


def _shipping_phone_bottom(slide) -> int:
    phone = _sticker_text_shapes(slide).get("phone")
    if phone is not None:
        return int(phone.top + phone.height)
    return int(Inches(4.63))


def _add_barcode_label(slide, barcode_value, left, top, width, height):
    """Left-aligned bag code under the bars, same width as the bar image."""
    box = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    box.name = "bagBarcodeLabel"
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = barcode_value
    p.alignment = PP_ALIGN.LEFT
    p.font.name = "Lato"
    p.font.size = Pt(BARCODE_LABEL_FONT_PT)
    try:
        box.fill.background()
        box.line.fill.background()
    except Exception:
        pass
    return box


def add_code128_barcode(slide, prs, barcode_value):
    left, barcode_width, top, barcode_height = _compute_barcode_rect(slide)
    text_h = int(BARCODE_TEXT_HEIGHT)
    bars_h = int(barcode_height) - text_h
    if bars_h < int(Inches(0.5)):
        text_h = int(barcode_height * 0.14)
        bars_h = int(barcode_height) - text_h

    bars_buf = _render_barcode_bars_png(barcode_value, barcode_width, bars_h)
    slide.shapes.add_picture(
        bars_buf,
        int(left),
        int(top),
        width=int(barcode_width),
        height=int(bars_h),
    )
    _add_barcode_label(
        slide,
        barcode_value,
        left,
        top + bars_h,
        barcode_width,
        text_h,
    )


def upsert_bag_to_airtable(db, shipping_info):
    dish_ids = [
        str(d["dishBarcode"])
        for d in shipping_info["Dishes"]
        if isinstance(d, dict) and d.get("dishBarcode")
    ]
    try:
        db.upsert_bag_record(
            bag_barcode=shipping_info["Bag Barcode"],
            dish_record_ids=dish_ids,
            ice_pack_required=shipping_info["Ice Pack Required"],
            shipping_name=shipping_info.get("Shipping Name"),
            zone=shipping_info.get("Zone Number"),
            household_members=shipping_info.get("Household Members"),
        )
    except Exception as e:
        logger.warning(f"Could not write bag {shipping_info['Bag Barcode']} to Airtable: {e}")


def create_shipping_stickers_barcode_ppt(db, shipping_list, template_path=None):
    if template_path is None:
        template_path = BASE_DIR / "template" / "Shipping_Sticker_Template_v3.pptx"

    prs = Presentation(str(template_path))

    if len(prs.slides) == 0:
        raise PPTGenerationError("Template presentation has no slides")

    template_slide = prs.slides[0]
    total_stickers = 0

    for shipping_info in shipping_list:
        stickers_needed = shipping_info["Stickers Needed"]

        for _ in range(stickers_needed):
            slide = copy_slide(template_slide, prs)
            populate_sticker(slide, shipping_info, shipping_info.get("Ice Pack Required", False))
            add_code128_barcode(slide, prs, shipping_info["Bag Barcode"])
            total_stickers += 1

        upsert_bag_to_airtable(db, shipping_info)

    # Remove original template slide.
    r_id = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(r_id)
    del prs.slides._sldIdLst[0]

    logger.info(f"Generated {total_stickers} shipping stickers with barcode")

    output = BytesIO()
    prs.save(output)
    output.seek(0)

    return output


def generate_shipping_stickers_barcode(db):
    template_path = BASE_DIR / "template" / "Shipping_Sticker_Template_v3.pptx"

    shipping_list = process_order_data(db)

    if not shipping_list:
        raise AirTableError("No shipping records to process")

    for info in shipping_list:
        dish_record_ids = [
            str(d.get("dishBarcode", ""))
            for d in info.get("Dishes", [])
            if str(d.get("dishBarcode", "")).startswith("rec")
        ]
        try:
            db.upsert_bag_record(
                bag_barcode=info["Bag Barcode"],
                dish_record_ids=dish_record_ids,
                ice_pack_required=info.get("Ice Pack Required", False),
                shipping_name=info.get("Shipping Name"),
                zone=info.get("Zone Number"),
                household_members=info.get("Household Members"),
            )
        except Exception as e:
            logger.warning(f"Could not write bag {info['Bag Barcode']} to Airtable: {e}")

    ppt_file = create_shipping_stickers_barcode_ppt(
        db,
        shipping_list,
        template_path=template_path,
    )

    return ppt_file, shipping_list


if __name__ == "__main__":
    try:
        db = new_database_access()

        ppt_file, shipping_list = generate_shipping_stickers_barcode(db)

        output_path = BASE_DIR / f"shipping_stickers_barcode_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"

        with open(output_path, "wb") as f:
            f.write(ppt_file.getvalue())

        logger.info(f"Generated {len(shipping_list)} unique bag records")
        logger.info(f"Saved PPT: {output_path}")
        logger.info("Bag rows written to Airtable Bag Tracking table")

    except AirTableError as e:
        logger.critical(f"Airtable error: {str(e)}")
    except Exception as e:
        logger.critical(f"Unexpected error: {str(e)}")