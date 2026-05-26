# generate_bag_stickers_barcode.py

import os
import json
import copy
import re
from datetime import datetime

import pandas as pd
from pptx import Presentation
from pptx.util import Inches

import sys
from pathlib import Path

BASE_DIR = Path(__file__).resolve().parents[2]  # Streamlit-App

sys.path.append(str(BASE_DIR))

from src.data.store_access import new_database_access
from src.stickers.dish_barcode_ids import dish_barcode_from_open_order_fields
from src.stickers.shipping_sticker_generator_v3 import (
    populate_sticker,
    copy_slide,
    add_code128_barcode,
    format_phone,
)

DEFAULT_BAG_TEMPLATE = BASE_DIR / "template" / "Shipping_Sticker_Template.pptx"

ICE_PACK_TAG = "Ice Pack"


def unwrap(value, default=""):
    """Airtable lookup fields sometimes return list values."""
    if isinstance(value, list):
        return value[0] if value else default
    if pd.isna(value):
        return default
    return value


def normalize_text(value):
    value = str(unwrap(value, "")).strip().lower()
    value = re.sub(r"\s+", " ", value)
    return value


def make_bag_group_key(row):
    """
    One bag = same delivery date + same shipping recipient/address + same zone.
    """
    parts = [
        normalize_text(row.get("Delivery Date")),
        normalize_text(row.get("Shipping Name")),
        normalize_text(row.get("Shipping Address 1")),
        normalize_text(row.get("Shipping Address 2")),
        normalize_text(row.get("Shipping Postal Code")),
        normalize_text(row.get("ZONE_NUMBER")),
    ]
    return "|".join(parts)


def customization_tags_from_row(row):
    """
    Merge tag values from any column whose name contains 'Customization Tags'
    (Airtable lookup / rollup names differ from the short label).
    """
    parts = []
    for key, val in row.items():
        if "Customization Tags" not in str(key):
            continue
        if isinstance(val, list):
            parts.extend(str(x) for x in val if str(x).strip())
        elif val is not None and not (isinstance(val, float) and pd.isna(val)) and str(val).strip():
            parts.append(str(val))
    return " ".join(parts)


def make_bag_barcode(row, index):
    """
    Stable readable bag barcode.
    Example: BAG-20260420-Z1-0001
    """
    delivery_date = str(unwrap(row.get("Delivery Date"), "")).replace("-", "")
    zone = str(unwrap(row.get("ZONE_NUMBER"), "X")).replace("Zone", "").strip()

    if not delivery_date:
        delivery_date = datetime.now().strftime("%Y%m%d")

    return f"BAG-{delivery_date}-Z{zone}-{index + 1:04d}"


def get_open_orders_df(db):
    """
    Pull open orders from Airtable.

    Important:
    This assumes your open orders view contains the same fields used by one pager:
    Delivery Date, Meal Sticker, Meal Portion, Customer Name, shipping fields,
    # of Parts, Zone Number, Shipping Name, and tags/customization info.
    """
    records = db.get_all_open_orders(view="viwDpTtU0qaT9NcvG")  

    rows = []
    for record in records:
        fields = record.get("fields", {})
        rows.append(fields)

    df = pd.DataFrame(rows).fillna("")
    df["Customization Tags"] = df.apply(customization_tags_from_row, axis=1)
    return df


def prepare_bag_dataframe(db):
    """
    Aggregate open orders into bag-level records.
    """
    df = get_open_orders_df(db)

    # Normalize possible Airtable/list fields.
    required_columns = [
        "Delivery Date",
        "Meal Sticker",
        "Meal Portion",
        "Customer Name",
        "Shipping Name",
        "Shipping Address 1",
        "Shipping Address 2",
        "Shipping City",
        "Shipping Province",
        "Shipping Postal Code",
        "Shipping Phone",
        "Zone Number (from Delivery Zone)",
        "# of Parts",
        "#",
        "Portion Result (in ClientServings)",
        "Customization Tags",
    ]

    for col in required_columns:
        if col not in df.columns:
            df[col] = ""

    df["Delivery Date"] = df["Delivery Date"].apply(unwrap)
    df["Meal Sticker"] = df["Meal Sticker"].apply(unwrap)
    df["Meal Portion"] = df["Meal Portion"].apply(unwrap)
    df["Customer Name"] = df["Customer Name"].apply(unwrap)
    df["Shipping Name"] = df["Shipping Name"].apply(unwrap)
    df["Shipping Address 1"] = df["Shipping Address 1"].apply(unwrap)
    df["Shipping Address 2"] = df["Shipping Address 2"].apply(unwrap)
    df["Shipping City"] = df["Shipping City"].apply(unwrap)
    df["Shipping Province"] = df["Shipping Province"].apply(unwrap)
    df["Shipping Postal Code"] = df["Shipping Postal Code"].apply(unwrap)
    df["Shipping Phone"] = df["Shipping Phone"].apply(unwrap) if "Shipping Phone" in df.columns else ""
    df["ZONE_NUMBER"] = df["Zone Number (from Delivery Zone)"].apply(unwrap)
    df["#"] = df["#"].apply(lambda x: unwrap(x, ""))

    # Build dish display text.
    df["dish_display"] = df.apply(
        lambda row: f"{row['Customer Name']} - {row['Meal Portion']} - {row['Meal Sticker']}",
        axis=1,
    )

    # Multi-part: one Client Serving row, but N physical stickers → N checklist lines (same dish #).
    expanded_rows = []

    for _, row in df.iterrows():
        parts = row.get("# of Parts", 1)
        parts = unwrap(parts, 1)

        try:
            parts = int(parts)
        except Exception:
            parts = 1

        if parts <= 1:
            new_row = row.copy()
            new_row["dish_display"] = row["dish_display"]
            expanded_rows.append(new_row)
        else:
            for part_num in range(1, parts + 1):
                new_row = row.copy()
                new_row["dish_display"] = f"{row['dish_display']} - PART {part_num}/{parts}"
                expanded_rows.append(new_row)

    df_expanded = pd.DataFrame(expanded_rows)

    # Bag grouping.
    df_expanded["BAG_GROUP_KEY"] = df_expanded.apply(make_bag_group_key, axis=1)

    bag_rows = []
    for bag_index, (_, group) in enumerate(df_expanded.groupby("BAG_GROUP_KEY", sort=False)):
        r0 = group.iloc[0]
        dishes = []
        household = set()
        tag_chunks = []

        for _, r in group.iterrows():
            cn = str(unwrap(r.get("Customer Name", "")))
            if cn.strip():
                household.add(cn)
            disp = str(r.get("dish_display", "")).strip()
            dishes.append(
                {
                    "dishBarcode": dish_barcode_from_open_order_fields(r.to_dict()),
                    "customerName": cn,
                    "mealPortion": str(unwrap(r.get("Meal Portion", ""))),
                    "mealSticker": str(unwrap(r.get("Meal Sticker", ""))),
                    "displayText": disp,
                    "status": "unscanned",
                }
            )
            tag_chunks.append(str(unwrap(r.get("Customization Tags", ""))))

        tags_merged = " ".join(tag_chunks)
        dish_list = "\n\n".join(d["displayText"] for d in dishes if d["displayText"])

        row_out = {
            "Delivery Date": unwrap(r0["Delivery Date"]),
            "Shipping Name": unwrap(r0["Shipping Name"]),
            "Shipping Address 1": unwrap(r0["Shipping Address 1"]),
            "Shipping Address 2": unwrap(r0["Shipping Address 2"]),
            "Shipping City": unwrap(r0["Shipping City"]),
            "Shipping Province": unwrap(r0["Shipping Province"]),
            "Shipping Postal Code": unwrap(r0["Shipping Postal Code"]),
            "Shipping Phone": format_phone(unwrap(r0.get("Shipping Phone", ""))),
            "ZONE_NUMBER": unwrap(r0["ZONE_NUMBER"]),
            "HOUSEHOLD_MEMBERS": "\n".join(sorted(household)),
            "DISH_LIST": dish_list,
            "TAGS_LIST": tags_merged,
            "DISH_OBJECTS": dishes,
            "TOTAL_ITEMS": len(dishes),
            "ICE_PACK_REQUIRED": ICE_PACK_TAG.lower() in tags_merged.lower(),
            "BAG_BARCODE": None,
        }
        row_out["BAG_BARCODE"] = make_bag_barcode(row_out, bag_index)
        row_out["line_shippingName"] = row_out["Shipping Name"]
        row_out["line_address"] = (
            f"{row_out['Shipping Address 1']}, {row_out['Shipping Address 2']}"
            if row_out["Shipping Address 2"]
            else row_out["Shipping Address 1"]
        )
        row_out["line_city"] = (
            f"{row_out['Shipping City']}, {row_out['Shipping Province']} {row_out['Shipping Postal Code']}"
        )
        row_out["line_zone"] = f"Zone {row_out['ZONE_NUMBER']}"
        row_out["line_totalItems"] = f"Total {row_out['TOTAL_ITEMS']} dish(es)"
        row_out["line_icePack"] = "Ice Pack Required" if row_out["ICE_PACK_REQUIRED"] else ""
        row_out["line_household"] = row_out["HOUSEHOLD_MEMBERS"]
        row_out["line_dishes"] = row_out["DISH_LIST"]
        bag_rows.append(row_out)

    return pd.DataFrame(bag_rows)


def copy_slide_with_images(source_slide, target_prs):
    target_slide = target_prs.slides.add_slide(source_slide.slide_layout)

    for shape in source_slide.shapes:
        try:
            new_element = copy.deepcopy(shape.element)
            target_slide.shapes._spTree.insert_element_before(new_element, "p:extLst")
        except Exception as e:
            print(f"Error copying shape: {e}")

    return target_slide


def _uses_one_pager_template(template_path):
    return "One_Pager" in str(template_path)


def _row_to_shipping_info(row):
    zone = str(row.get("ZONE_NUMBER", "")).replace("Zone", "").replace("ZONE", "").strip()
    return {
        "Shipping Name": row["Shipping Name"],
        "Shipping Address 1": row["Shipping Address 1"],
        "Shipping Address 2": row["Shipping Address 2"],
        "Shipping City": row["Shipping City"],
        "Shipping Province": row["Shipping Province"],
        "Shipping Postal Code": row["Shipping Postal Code"],
        "Shipping Phone": row.get("Shipping Phone", ""),
        "Zone Number": zone,
    }


def populate_bag_slide(slide, row, template_path):
    if _uses_one_pager_template(template_path):
        populate_slide(slide, row)
    else:
        populate_sticker(
            slide,
            _row_to_shipping_info(row),
            ice_pack_required=bool(row.get("ICE_PACK_REQUIRED", False)),
        )


def populate_slide(slide, row):
    """
    Fill placeholders by matching PPT shape names to dataframe columns.
    Example PPT shape names:
    line_shippingName, line_address, line_city, line_zone,
    line_totalItems, line_icePack, line_household, line_dishes, BAG_BARCODE
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        key = shape.name

        if key in row:
            value = row.get(key, "")
        elif key == "BAG_BARCODE_TEXT":
            value = row.get("BAG_BARCODE", "")
        else:
            continue

        text_frame = shape.text_frame
        if text_frame.paragraphs and text_frame.paragraphs[0].runs:
            text_frame.paragraphs[0].runs[0].text = str(value)
        else:
            text_frame.text = str(value)


def generate_bag_stickers_barcode(
    db,
    template_path=None,
    export_mapping_path=None,
):
    if template_path is None:
        template_path = DEFAULT_BAG_TEMPLATE
    template_path = Path(template_path)

    df = prepare_bag_dataframe(db)

    template_path = Path(template_path)
    use_shipping_layout = not _uses_one_pager_template(template_path)

    prs = Presentation(str(template_path))
    template_slide = prs.slides[0]

    for _, row in df.iterrows():
        if use_shipping_layout:
            new_slide = copy_slide(template_slide, prs)
        else:
            new_slide = copy_slide_with_images(template_slide, prs)
        populate_bag_slide(new_slide, row, template_path)
        add_code128_barcode(new_slide, prs, row["BAG_BARCODE"])

    # Remove template slide.
    r_id = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(r_id)
    del prs.slides._sldIdLst[0]

    mapping = []
    for _, row in df.iterrows():
        dishes = row.get("DISH_OBJECTS") or [
            dish.strip()
            for dish in str(row["DISH_LIST"]).split("\n\n")
            if dish.strip()
        ]
        mapping.append(
            {
                "bagBarcode": row["BAG_BARCODE"],
                "deliveryDate": row["Delivery Date"],
                "shippingName": row["Shipping Name"],
                "address1": row["Shipping Address 1"],
                "address2": row["Shipping Address 2"],
                "city": row["Shipping City"],
                "province": row["Shipping Province"],
                "postalCode": row["Shipping Postal Code"],
                "zone": row["ZONE_NUMBER"],
                "householdMembers": row["HOUSEHOLD_MEMBERS"].split("\n"),
                "icePackRequired": bool(row["ICE_PACK_REQUIRED"]),
                "totalItems": int(row["TOTAL_ITEMS"]),
                "dishes": dishes,
            }
        )

    if export_mapping_path:
        with open(export_mapping_path, "w") as f:
            json.dump(mapping, f, indent=2)

    # Write each bag to Airtable Bag Tracking (source of truth for plating website).
    for entry in mapping:
        dish_record_ids = list(
            dict.fromkeys(
                str(d["dishBarcode"])
                for d in entry["dishes"]
                if isinstance(d, dict) and str(d.get("dishBarcode", "")).startswith("rec")
            )
        )
        try:
            db.upsert_bag_record(
                bag_barcode=entry["bagBarcode"],
                dish_record_ids=dish_record_ids,
                ice_pack_required=entry["icePackRequired"],
                shipping_name=entry.get("shippingName"),
                zone=entry.get("zone"),
                household_members=entry.get("householdMembers"),
            )
        except Exception as e:
            print(f"⚠️ Could not write bag {entry['bagBarcode']} to Airtable: {e}")

    return prs, df


if __name__ == "__main__":
    db = new_database_access()

    prs, df = generate_bag_stickers_barcode(
        db,
        template_path=str(DEFAULT_BAG_TEMPLATE),
        export_mapping_path=None,
    )

    output_path = BASE_DIR / f"bag_stickers_barcode_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    prs.save(output_path)

    print(f"Generated {len(df)} bag barcode stickers")
    print(f"Saved PPT: {output_path}")
    ice_count = int(df["ICE_PACK_REQUIRED"].sum()) if "ICE_PACK_REQUIRED" in df.columns else 0
    if ice_count:
        print(f"Ice-pack bags ({ice_count}): snowflake on name line (B&W friendly)")
    print(f"Saved mapping: {BASE_DIR / 'bag_barcode_mapping.json'}")