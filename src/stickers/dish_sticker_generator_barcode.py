# Standard library imports
import os
import json
import glob
from io import BytesIO
from datetime import datetime, timedelta
# Third-party imports
import barcode
import time
from barcode.writer import ImageWriter
import brother_ql
from brother_ql.conversion import convert
from brother_ql.backends.helpers import send
from dotenv import load_dotenv
import pandas as pd
from PIL import Image
import pytz
from pptx import Presentation
from pptx.util import Inches
from pyairtable.api.table import Table
import streamlit as st
from src.data.store_access import new_database_access
import gc

# Local imports
import copy
est = pytz.timezone('US/Eastern')
current_date_time = datetime.now(est).strftime("%Y%m%d_%H%M")

# BASE_ID = "appkvckTQ86vhjSC0"  # 041025 
BASE_ID = "appEe646yuQexwHJo" # production
TABLE_ID = "tblVwpvUmsTS2Se51"  # ClientServings
VIEW_ID = "viw5hROs9I9vV0YEq" # Sorted (For Dish Sticker)

# Set the time zone to Eastern Standard Time (EST)
est = pytz.timezone('US/Eastern')

class AirTableError(Exception):
    """Custom exception for AirTable related errors"""
    pass

class PPTGenerationError(Exception):
    """Custom exception for PowerPoint generation errors"""
    pass


def resize_image(image_path, target_width, target_height):
    """
    Resize the image to the specified width and height.

    Args:
        image_path (str): Path to the image file.
        target_width (int): Target width for the resized image.
        target_height (int): Target height for the resized image.
    """
    try:
        with Image.open(image_path) as img:
            img = img.resize((target_width, target_height), Image.Resampling.LANCZOS)
            img.save(image_path)
        print(f"Successfully resized {image_path} to {target_width}x{target_height}")
    except Exception as e:
        print(f"Error resizing {image_path}: {str(e)}")

def ensure_file_saved(file_path, max_retries=5, delay=0.1):
    """
    Ensure a file is completely saved before proceeding.
    
    Args:
        file_path (str): Path to the file to check
        max_retries (int): Maximum number of retry attempts
        delay (float): Delay between retries in seconds
    
    Returns:
        bool: True if file is saved and accessible, False otherwise
    """
    file_path = file_path
    for attempt in range(max_retries):
        try:
            # Check if file exists and is readable
            if os.path.exists(file_path) and os.access(file_path, os.R_OK):
                # Try to open the file to ensure it's not locked
                with open(file_path, 'rb') as f:
                    f.read(1)  # Try to read at least 1 byte
                return True
        except (OSError, IOError):
            pass
        
        time.sleep(delay)
    
    return False
    
def sort_key(filename):
    """
    Extract the numeric part from the filename for sorting.

    Args:
        filename (str): The filename to extract the number from.

    Returns:
        int: The extracted number for sorting.
    """
    last_part = filename.split('-')[-1]
    number = int(last_part.split('.')[0])
    return number

def print_stickers():
    """
    Print the generated dish stickers using Brother QL printer.
    """
    current_file_path = os.path.abspath(__file__)
    
    # Set up printer configuration
    model = 'QL-810W'
    
    # Configure source folder and printer settings
    source_folder = "/Users/clairegoldwitz/Desktop/image_test/test.jpg"
    backend = 'pyusb'
    identifier = 'usb://0x04f9:0x209c'
    
    # Alternative configuration for network printing
    # source_folder = os.path.join(os.path.dirname(current_file_path), 'assets', 'images')
    # backend = 'network'
    # identifier = 'tcp://192.168.4.30'
    
    # Load all JPG and PNG files in order
    try:
        label_files = glob.glob(f'{source_folder}/*.jpg') + glob.glob(f'{source_folder}/*.png')
        label_files.sort(key=sort_key, reverse=True)
        print(f"Found {len(label_files)} image files")
    except Exception as e:
        print(f"Error loading image files: {str(e)}")
        return
    
    if not label_files:
        print("No image files found in the directory.")
        return
    
    for i, label_file in enumerate(label_files):
        # print(f'LABEL FILE: {label_file} | {i+1} of {len(label_files)}')
        # Create the label maker object
        try:
            qlr = brother_ql.BrotherQLRaster(model)
            qlr.exception_on_warning = True
            print("Successfully created label maker object")
        except Exception as e:
            print(f"Error creating label maker object: {str(e)}")
            return
        
        img = Image.open(label_file)
        # img = img.resize((3800, 2512))
        
        try:
            instructions = convert(
                qlr=qlr,
                images=[img],
                label='62',
                cut=True,
                rotate='90'
            )
            print(f"Successfully converted label {i+1}")
        except Exception as e:
            print(f"Error converting label {i+1}: {str(e)}")
            return
        
        try:
            send(instructions=instructions, printer_identifier=identifier, backend_identifier=backend)
            print(f"Successfully sent label {i+1} to printer")
        except Exception as e:
            print(f"Error sending label {i+1} to printer: {e}")
            return
            
        # Delete previous file:
        try:
            if i > 0 and os.path.exists(label_files[i - 1]):
                print(f"Deleting previous file: {label_files[i - 1]}")
                os.remove(label_files[i - 1])
        except Exception as e:
            print(f"Error deleting previous file: {e}. Check permissions.")
            
    try:
        if os.path.exists(label_files[-1]):
            print(f"Deleting previous file: {label_files[-1]}")
            os.remove(label_files[-1])
    except Exception as e:
        print(f"Error deleting previous file: {e}. Check permissions.")
            
    print("\nPrinting process completed.")

# Process Data 
def generate_sticker_df(df):
    def add_days(date_str, days = 3):  # by default: add 3 days based on EST
        date = datetime.now(est).strptime(date_str, "%Y-%m-%d")
        end_date = date + timedelta(days)
        return end_date.strftime('%m/%d/%Y')
    
    df['Best Before'] = df['Delivery Date'].apply(lambda x: add_days(x[0] if isinstance(x, list) else x, 3))
    df['best_before'] = "Best before: " + df['Best Before']

    meal_info = df['Meal Sticker (from Linked OrderItem)'].apply(lambda x: x[0] if isinstance(x, list) else x)
    df[['bowl_name', 'ingredients']] = meal_info.str.split(': ', n=1, expand=True)
    df['parts'] = df['# of Parts'].apply(lambda x: x[0] if isinstance(x, list) else x)

    # duplicate according to '# portions'
    # df_dish = df_dish.loc[df_dish.index.repeat(df_dish['# portions'])]
    # df_dish = df_dish.reset_index(drop=True)
    
    return df

"""### Draw Slides"""

def insert_background(slide, img_path, prs):
    left = top = Inches(0)
    pic = slide.shapes.add_picture(img_path, left, top, 
                                width=prs.slide_width, height=prs.slide_height)
    # This moves it to the background
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)
    return slide

def copy_slide(template, prs):
    # Create a new slide with the same layout as the source slide
    new_slide = prs.slides.add_slide(template.slide_layout)

    # Copy the content from the source slide to the new slide
    for shape in template.shapes:
        new_shape = copy.deepcopy(shape)
        new_slide.shapes._spTree.insert_element_before(new_shape._element, 'p:extLst')

    return new_slide

def generate_dish_stickers_barcode(db, progress_placeholder=None, cancel_event=None, progress=None):
    """
    Generate dish stickers with barcode.

    Args:
        db: Database access object
        progress_placeholder: Optional Streamlit placeholder for progress updates (sync use)
        cancel_event: Optional threading.Event — when set, stop generating and return partial prs
        progress: Optional dict — status / slide_count / total_slides are written here for the
            caller to read from another thread (avoids st.* calls from a worker)
    """
    client_serving_df = read_client_serving(db)
    print(client_serving_df.head())
    df = generate_sticker_df(client_serving_df)
    # Load the presentation
    prs = Presentation('template/Dish_Sticker_Template_Barcode.pptx')

    # Get the first slide
    slide = prs.slides[0]

    # https://python-barcode.readthedocs.io/en/stable/getting-started.html#usage
    ITF = barcode.get_barcode_class('ITF')
    BARCODE_HEIGHT = 0.75
    margin = Inches(0.05)


    # Find the maximum number of parts across all rows
    max_parts = df['parts'].max()
    total_slides = df['parts'].sum()
    print(f"Max parts: {max_parts}, Total rows: {len(df)}, Total slides to generate: {total_slides}")
    if progress is not None:
        progress["total_slides"] = int(total_slides)
        progress["slide_count"] = 0

    slide_count = 0
    # Iterate row-by-row so all parts of the same client serving stay adjacent
    # (Part 1 immediately followed by Part 2, etc.)
    for idx, row in df.iterrows():
        parts_count = int(row['parts'])
        for part_num in range(parts_count):
            if cancel_event is not None and cancel_event.is_set():
                return prs
            slide_count += 1
            if slide_count % 20 == 0:
                gc.collect()

            status_msg = f"Slide {slide_count}/{total_slides} - {row['client_name']} ({row['#']})"
            print(f"  {status_msg}")

            if progress is not None:
                progress["status"] = status_msg
                progress["slide_count"] = slide_count

            # Update Streamlit progress if placeholder provided
            if progress_placeholder is not None:
                progress_placeholder.text(status_msg)
            parts = row['parts']

            # add one page copied from template
            new_slide = copy_slide(slide, prs)

            for shape in new_slide.shapes:
                if shape.has_text_frame:
                    key = shape.name
                    text_frame = shape.text_frame
                    if text_frame.paragraphs and text_frame.paragraphs[0].runs:
                        value = row.get(key, 'N/A')
                        if key == 'parts':
                            if parts > 1:
                                value = f"PART {part_num+1}/{parts}"
                            else:
                                value = ''
                            #print(f"DEBUG: key={key}, parts={parts}, value={value}")
                        text_frame.paragraphs[0].runs[0].text = str(value) if value is not None else 'N/A'

            # ITF barcode generator, directly to memory (no disk I/O)
            itf = ITF(row["#"], writer=ImageWriter(mode="RGBA"))
            barcode_buffer = BytesIO()
            itf.write(barcode_buffer, dict(quiet_zone=3, background=(255, 255, 255, 0),
                                           font_size=5, text_distance=2.5,
                                           module_width=0.3))
            barcode_buffer.seek(0)

            id_barcode = Image.open(barcode_buffer)
            barcode_size = ((Inches(id_barcode.size[0] / id_barcode.size[1] * BARCODE_HEIGHT)), Inches(BARCODE_HEIGHT))

            # Reset buffer position for pptx to read
            barcode_buffer.seek(0)
            new_slide.shapes.add_picture(barcode_buffer, prs.slide_width - margin - barcode_size[0],
                                         prs.slide_height - margin - barcode_size[1],
                                         width=barcode_size[0], height=barcode_size[1])
            id_barcode.close()
    return prs

def read_client_serving(db):
    fields_to_return = ['#', 'Customer Name', 'Meal Sticker (from Linked OrderItem)', 'Meal Portion (from Linked OrderItem)', 'Delivery Date', 'Position Id', 'Dish ID (from Linked OrderItem)','Delivery Zone (from Linked OrderItem)','# of Parts','Meal Type from Profile (from Linked OrderItem)']
    try:
        # Initialize table
        records = db.get_clientservings_data(view=VIEW_ID)
        
        if not records:
            raise AirTableError("No records found in the specified view.")
        
        # Convert records to DataFrame
        df = pd.DataFrame(records)
        
        # Extract fields from the records
        if 'fields' in df.columns:
            fields_df = pd.json_normalize(df['fields'])
            # Combine the records with their fields
            df_flat = pd.concat([df.drop('fields', axis=1), fields_df], axis=1)
        else:
            df_flat = df
        
        if 'Dish ID (from Linked OrderItem)' in df_flat.columns and 'Position Id' in df_flat.columns:
            df_flat['Dish ID (from Linked OrderItem)'] = df_flat['Dish ID (from Linked OrderItem)'].apply(
                lambda x: int(x[0]) if isinstance(x, list) and len(x) > 0 and str(x[0]).isdigit() 
                else (int(x) if isinstance(x, (int, float)) else x)
            )
        df_flat = df_flat[fields_to_return]

        add_ons = db.get_all_add_ons()
        breakfast_dishes = db.get_all_breakfast_dishes()
        print(f"Add-ons: {add_ons}")
        print(f"Breakfast dishes: {breakfast_dishes}")

        df_flat = df_flat.fillna('N/A')
    
        # Generate Requested Columns
        # barcode ID: got a float - change to string
        df_flat["#"] = df_flat["#"].apply(lambda x: str(int(x)))

        # all linked item needs to be extracted from a list
        df_flat['client_name'] = df_flat['Customer Name'].apply(lambda x: x[0] if isinstance(x, list) else x)

        # zone
        df_flat['zone'] = df_flat['Delivery Zone (from Linked OrderItem)'].apply(lambda x: 'Zone ' + str(x[0]) if isinstance(x, list) else 'Zone ' + str(x))

        #mark record yes in "add-on" if dish id is in add-ons
        df_flat['add-on'] = df_flat['Dish ID (from Linked OrderItem)'].apply(lambda x: 'yes' if x in add_ons else 'no')
        
        # Derive meal from 'Meal Type from Profile (from Linked OrderItem)' (linked field arrives as a list)
        df_flat['meal'] = df_flat['Meal Type from Profile (from Linked OrderItem)'].apply(
            lambda x: (x[0] if isinstance(x, list) and x else x) or ''
        )

        # Sort order: lunch/dinner → breakfast → lunch&dinner add-ons → snacks → snacks add-ons
        # A dish flagged as Breakfast in the weekly menu is ALWAYS treated as breakfast,
        # even if the customer's order meal portion says lunch/dinner.
        def _sort_rank(row):
            meal = str(row['meal']).strip().lower()
            dish_id = row['Dish ID (from Linked OrderItem)']
            is_addon = row['add-on'] == 'yes'
            is_snack = 'snack' in meal
            is_breakfast = 'breakfast' in meal or dish_id in breakfast_dishes
            if is_snack:
                return 4 if is_addon else 3
            if is_addon:
                return 2  # lunch/dinner add-ons
            if is_breakfast:
                return 1
            return 0  # lunch/dinner non-add-on

        # Sort: tier first (dish-level), then group by Dish ID, then Position Id within dish
        df_flat['_sort_rank'] = df_flat.apply(_sort_rank, axis=1)
        df_flat = df_flat.sort_values(
            by=['_sort_rank', 'Dish ID (from Linked OrderItem)', 'Position Id'],
            ascending=[True, True, True]
        ).drop(columns=['_sort_rank'])
        #save to csv
        # df_flat.to_csv(f'{current_date_time}_dish_sticker_barcode.csv', index=False)
        return df_flat
        
    except Exception as e:
        raise AirTableError(f"Error fetching data from Airtable: {str(e)}")

if __name__ == "__main__":
    db = new_database_access()

    prs = generate_dish_stickers_barcode(db)
    prs.save(f'{current_date_time}_dish_sticker_barcode.pptx')

    #print_stickers()