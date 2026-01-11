import streamlit as st
import pandas as pd
from pptx import Presentation
import copy
from datetime import datetime, timedelta
from pptx.util import Inches, Pt
from pyairtable.api.table import Table
from pyairtable.formulas import match
from dotenv import load_dotenv
from functools import cache
import math
import os
import xml.etree.ElementTree as ET
from lxml import etree
from pptx.oxml import parse_xml
from store_access import new_database_access
ICE_PACK_TAG = 'Ice Pack'

def get_open_orders(db):
    open_orders_fields = {
        'CLIENT': 'fldjEgeRh2bGxajXT',            
        'DISH_STICKER': 'fldYZYDRjScz6ig5a',       
        'DELIVERY_DAY': 'flddRJziNBdEtrpmG', 
        'PORTIONS': 'fldE0fWRfUnoHznqC',
        'SHIPPING_ADDRESS_1': 'flddnU6Y02iIHp16G',
        'SHIPPING_CITY': 'fldBIgt7ce5fEYLDG',
        'SHIPPING_PROVINCE': 'fldGeLuYoGKiw227Y',
        'SHIPPING_COUNTRY': 'fldnP3H74kAXKaIhN',
        'SHIPPING_POSTAL_CODE': 'fldwqwf7WSbiP0hJg',
        'SHIPPING_ADDRESS_2': 'fldRUUiFRRYQ52k0W',
        'CUSTOMER_NAME': 'fldWYJStYSpX72pG3',
        'PARTS': 'fldebKYuTeuQfauil',
        'ZONE_NUMBER': 'fldbL18Ixas6ong0j'   
    }
    data = db.get_all_open_orders(view='viwuVy9aN2LLZrcPF')
    # Create DataFrame and map column names
    df = pd.DataFrame([record['fields'] for record in data])
    column_mapping = {v.replace('fld', ''): k for k, v in open_orders_fields.items()}
    df.rename(columns=column_mapping, inplace=True)
    
    return df

def get_clients_info(db):
    clients_fields = {
        'IDENTIFIER': 'fldDVCUtcmEv5ZkEv',              
        'KCAL': 'fldgaURVgPSD6UFVU',
        'CARBS': 'fldRohbKpwyIyyuhe',
        'PROTEIN': 'fldLpIgHNZFxyZ1PI',
        'FAT': 'fldMSEudjCwovSD8A',
        'FIBER': 'fld1TirSqBsA1GqK5',
        'CLIENT_FNAME': 'fldGBn8BpGwFIHxgi',
        'CLIENT_LNAME': 'fldIq6giA1dDcut8T',
        'TAGS':'fldKdacQ9GMB070cD'
    }
    data = db.get_all_clients()  # Replace with actual view ID 
    # Create DataFrame and map column names
    df = pd.DataFrame([{**record['fields'], 'id': record['id']} for record in data])
    column_mapping = {v.replace('fld', ''): k for k, v in clients_fields.items()}
    df.rename(columns=column_mapping, inplace=True)
    # print('Data from clients table:', df)
    return df

def process_data(db):
    # Get data from both tables
    df_orders = get_open_orders(db)
    df_clients = get_clients_info(db)
    
    # Clean column names and values
    df_orders.columns = [col.strip() for col in df_orders.columns]
    df_clients.columns = [col.strip() for col in df_clients.columns]
    required_cols = [
        'Delivery Date', 'Meal Sticker', 'Meal Portion', 'To_Match_Client_Nutrition',
        'Shipping Address 1', 'Shipping Address 2', 'Shipping City', 
        'Shipping Province', 'Shipping Postal Code', 'Customer Name','# of Parts','Zone Number (from Delivery Zone)'
    ]

    # Ensure all required columns exist (missing ones will be created and filled with "")
    df_orders = df_orders.reindex(columns=required_cols, fill_value="")

    # Drop rows where 'To_Match_Client_Nutrition' is missing/blank
    df_orders = df_orders[df_orders['To_Match_Client_Nutrition'] != ""]
    
    # Convert any list values to strings in all relevant columns
    # For client identifier
    df_orders["To_Match_Client_Nutrition"] = df_orders["To_Match_Client_Nutrition"].apply(
        lambda x: x[0] if isinstance(x, list) and len(x) > 0 else 
                (str(x) if not isinstance(x, list) else "Unknown")
    )
    
    # For meal sticker
    df_orders["Meal Sticker"] = df_orders.apply(
        lambda row: (row["Meal Sticker"][0] if isinstance(row["Meal Sticker"], list) and len(row["Meal Sticker"]) > 0 else 
                    (str(row["Meal Sticker"]) if not isinstance(row["Meal Sticker"], list) else "Unknown")) 
                    + ' - ' + str(row.get("Meal Portion", "Unknown")),
        axis=1
    )
    
    # For order type
    df_orders["Meal Portion"] = df_orders["Meal Portion"].apply(
        lambda x: x[0] if isinstance(x, list) and len(x) > 0 else 
                (str(x).replace(' Subscriptions', '') if not isinstance(x, list) else "Unknown")
    )
    
    # Determine portions based on order type
    df_orders['Portions'] = df_orders.apply(
        lambda row: 0.5 if row['Meal Portion'] == 'Breakfast' 
                else (0.25 if row['Meal Portion'] == 'Snack' else 1),
        axis=1
    )
    
    # Calculate portion strings with appropriate line breaks
    standard_str = "Fish Taco Bowl with Grilled Chicken, Shredded Red and Green Cabbage, Roasted"
    unit_len = len(standard_str)
    
    df_orders['portion_str'] = df_orders.apply(
        lambda row: '[ 1 ] ' + 
        '\n' * (math.ceil(len(row['Meal Sticker'])/unit_len) - 1), 
        axis=1
    )
    
    # Rename client column for consistent merging
    df_orders.rename(columns={'To_Match_Client_Nutrition': 'CLIENT'}, inplace=True)
    df_orders.rename(columns={'Zone Number (from Delivery Zone)': 'ZONE_NUMBER'}, inplace=True)
    df_clients.rename(columns={'id': 'CLIENT'}, inplace=True)
    
    # Ensure client names are cleaned
    df_orders.CLIENT = df_orders.CLIENT.apply(lambda x: str(x).strip() if isinstance(x, str) else x)
    df_clients.CLIENT = df_clients.CLIENT.apply(lambda x: str(x).strip() if isinstance(x, str) else x)

    df_orders.ZONE_NUMBER = df_orders.ZONE_NUMBER.apply(lambda x: str(x[0] if isinstance(x, list) else str(x).strip()))
    
    # Process nutrition information
    for _, row in df_clients.iterrows():
        if 'Customization Tags' in row and isinstance(row['Customization Tags'], list) and "No Nutrition Data in Sheet" in row['Customization Tags']:
            nutrition_line = ''
        elif 'goal_calories' in row and not pd.isna(row['goal_calories']) and row['goal_calories'] > 0:
            nutrition_line = row['identifier'].split('|')[1].strip() + ' Nutrition target per serving: '
            nutrition_line += f"{int(row['goal_calories'])} kcals | "
            if 'goal_carbs(g)' in row and not pd.isna(row['goal_carbs(g)']) and row['goal_carbs(g)'] > 0:
                nutrition_line += f"{int(row['goal_carbs(g)'])}g carbs, "
            if 'goal_protein(g)' in row and not pd.isna(row['goal_protein(g)']) and row['goal_protein(g)'] > 0:
                nutrition_line += f"{int(row['goal_protein(g)'])}g protein, "
            if 'goal_fat(g)' in row and not pd.isna(row['goal_fat(g)']) and row['goal_fat(g)'] > 0:
                nutrition_line += f"{int(row['goal_fat(g)'])}g fat, "
            if 'goal_fiber(g)' in row and not pd.isna(row['goal_fiber(g)']) and row['goal_fiber(g)'] > 0:
                nutrition_line += f"{int(row['goal_fiber(g)'])}g fiber"
        
        df_clients.loc[_, 'NUTRITION'] = nutrition_line
    
    # Group the nutrition info
    df_clients['NUTRITION'] = df_clients.groupby('CLIENT')['NUTRITION'].apply(
        lambda x: '\n'.join(filter(None, x))
    ).reset_index(drop=True)
    
    # Extract the relevant parts of the ID for grouping
    df_clients['group_key'] = df_clients['identifier'].apply(lambda x: "".join(x.split("|")[::2]) if isinstance(x, str) else x)
    
    # Sort and consolidate nutrition lines
    meal_order = ['breakfast', 'lunch', 'dinner', 'snack']
    
    def sort_nutrition_lines(lines):
        def meal_sort_key(line):
            for i, meal in enumerate(meal_order):
                if line.lower().startswith(meal):
                    return i
            return len(meal_order)  # if not found, put it at the end
        sorted_lines = sorted(filter(None, lines), key=meal_sort_key)
        return '\n'.join(sorted_lines)
    
    consolidated_nutrition_map = df_clients.groupby('group_key')['NUTRITION'].apply(
        sort_nutrition_lines
    ).to_dict()
    
    # Map the consolidated nutrition back to each row
    df_clients['consolidated_nutrition'] = df_clients['group_key'].map(consolidated_nutrition_map)
    df_clients['NUTRITION'] = df_clients['consolidated_nutrition'].fillna('')
    
    # IMPROVEMENT: First merge the dataframes, then calculate page indices
    df_merge = df_orders.merge(df_clients, on='CLIENT', how='left')

    # Expand dishes with parts > 1 BEFORE page calculation
    # For each dish with parts > 1, create multiple rows with "PART x/y" suffix
    def expand_dish_parts(row):
        parts = row.get('# of Parts', 1)
        if pd.isna(parts) or parts == '' or parts is None:
            parts = 1
        parts = int(parts)

        if parts > 1:
            expanded_stickers = []
            expanded_portions = []
            for i in range(1, parts + 1):
                expanded_stickers.append(f"{row['Meal Sticker']} - PART {i}/{parts}")
                expanded_portions.append(row['portion_str'])
            return pd.Series({'Meal Sticker': expanded_stickers, 'portion_str': expanded_portions})
        else:
            return pd.Series({'Meal Sticker': [row['Meal Sticker']], 'portion_str': [row['portion_str']]})

    # Apply expansion and explode the lists
    expanded = df_merge.apply(expand_dish_parts, axis=1)
    df_merge['Meal Sticker'] = expanded['Meal Sticker']
    df_merge['portion_str'] = expanded['portion_str']
    df_merge = df_merge.explode(['Meal Sticker', 'portion_str']).reset_index(drop=True)

    # Calculate page indices AFTER expansion (so expanded rows count toward page size)
    page_size = 6
    df_merge['Index'] = df_merge.groupby(['Customer Name','Shipping Address 1']).cumcount()
    df_merge['page_number'] = df_merge['Index'] // page_size

    # Function to calculate total characters in meal stickers for a group
    def get_total_chars(group):
        return group['Meal Sticker'].str.len().sum()

    # Split pages with long meal sticker names
    df_merge['temp_group'] = df_merge.groupby(['Customer Name', 'Shipping Address 1', 'page_number']).ngroup()
    page_splits = []

    for group_id in df_merge['temp_group'].unique():
        group = df_merge[df_merge['temp_group'] == group_id].copy()
        if len(group) == 6 and get_total_chars(group) > 1000:  # If page is full and has long names
            # Split into two pages: 5 items + ice pack, and 1 item + ice pack
            first_page = group.iloc[:5].copy()
            second_page = group.iloc[5:].copy()

            # Update page numbers
            first_page['page_number'] = group['page_number'].iloc[0]
            second_page['page_number'] = group['page_number'].iloc[0] + 1

            # Update indices
            first_page['Index'] = range(len(first_page))
            second_page['Index'] = range(len(second_page))

            page_splits.append((first_page, second_page))
        else:
            page_splits.append((group, None))

    # Reconstruct dataframe with split pages
    new_dfs = []
    for first, second in page_splits:
        new_dfs.append(first)
        if second is not None:
            new_dfs.append(second)

    df_merge = pd.concat(new_dfs, ignore_index=True)
    df_merge = df_merge.drop('temp_group', axis=1)

    # Add ice pack rows for customers who have ice pack tag
    customers_with_ice_pack = df_merge[df_merge['Customization Tags'].apply(lambda x: isinstance(x, list) and ICE_PACK_TAG in x)]['Customer Name'].unique()

    ice_pack_rows = []
    for customer in customers_with_ice_pack:
        # Get all pages for this customer
        customer_pages = df_merge[df_merge['Customer Name'] == customer]
        max_page = int(customer_pages['page_number'].max())

        # Add ice pack to each page
        for page in range(max_page + 1):
            # Get items for this page
            page_items = customer_pages[customer_pages['page_number'] == page]
            if not page_items.empty:
                # Get the last row for this page to copy details
                customer_row = page_items.iloc[-1].copy()
                # Create ice pack row
                ice_pack_row = customer_row.copy()
                ice_pack_row['Meal Sticker'] = 'Ice Pack'
                ice_pack_row['Index'] = page_items['Index'].max() + 1  # Make it the next index after the last item
                ice_pack_row['page_number'] = page  # Keep it on the same page
                ice_pack_rows.append(ice_pack_row)

    if ice_pack_rows:
        df_merge = pd.concat([df_merge, pd.DataFrame(ice_pack_rows)], ignore_index=True)

    # Sort so ice pack rows are at the bottom of each group
    df_merge['is_ice_pack'] = df_merge['Meal Sticker'] == 'Ice Pack'
    df_merge = df_merge.sort_values(['Customer Name', 'Shipping Address 1', 'is_ice_pack'])

    # Remove the temporary sorting column
    df_merge = df_merge.drop('is_ice_pack', axis=1)

    # Simple groupby to create page-level information
    df_grouped = df_merge.groupby(['Customer Name', 'page_number']).agg({
        # Keep the first occurrence of these customer details per page
        'First_Name': 'first',
        'Last_Name': 'first',
        'Delivery Date': 'first',
        'Shipping Address 1': 'first',
        'Shipping Address 2': 'first',
        'Shipping City': 'first',
        'Shipping Province': 'first',
        'Shipping Postal Code': 'first',
        'NUTRITION': 'first',
        'ZONE_NUMBER': 'first',
        # Combine these for each page
        'portion_str': lambda x: '\n\n'.join(x),
        'Meal Sticker': lambda x: '\n\n'.join(x),
        
    }).reset_index()

    # Rename the aggregated columns to match expected names
    df_grouped.rename(columns={
        'portion_str': 'portion_list',
        'Meal Sticker': 'dish_list'
    }, inplace=True)

    # Create a unique client identifier
    df_grouped['CLIENT_UNIQUE_ID'] = df_grouped['First_Name'].fillna('') + ' ' + df_grouped['Last_Name'].fillna('')
    df_grouped['CLIENT_UNIQUE_ID'] = df_grouped['CLIENT_UNIQUE_ID'].apply(lambda x: x.strip())
    df_merge_grouped = df_grouped.sort_values(by=['CLIENT_UNIQUE_ID'])
    
    # Create a concatenated shipping address for household grouping
    address_components = ['Shipping Address 1', 'Shipping Address 2', 'Shipping Postal Code']
    has_address_fields = any(field in df_merge_grouped.columns for field in address_components)
    
    if has_address_fields:
        # Normalize and combine address components
        for field in address_components:
            if field in df_merge_grouped.columns:
                df_merge_grouped[field] = df_merge_grouped[field].fillna('').astype(str).str.lower().str.strip()
        
        # Create a combined address string for grouping
        df_merge_grouped['COMPLETE_ADDRESS'] = ''
        for field in address_components:
            if field in df_merge_grouped.columns:
                df_merge_grouped['COMPLETE_ADDRESS'] += df_merge_grouped[field] + '|'
        
        # Group orders by client and get the first address for each client
        client_addresses = df_merge_grouped.groupby('CLIENT_UNIQUE_ID').agg({
            'COMPLETE_ADDRESS': 'first'
        }).reset_index()
        
        # Create a dictionary to map clients to their addresses
        address_map = dict(zip(client_addresses['CLIENT_UNIQUE_ID'], client_addresses['COMPLETE_ADDRESS']))
        
        # Apply this mapping to get consistent addresses
        df_merge_grouped['HOUSEHOLD_GROUP'] = df_merge_grouped['CLIENT_UNIQUE_ID'].map(address_map)
        df_merge_grouped['HOUSEHOLD_GROUP'] = df_merge_grouped['HOUSEHOLD_GROUP'].fillna(df_merge_grouped['CLIENT_UNIQUE_ID'])
        
        print("Grouping households by shipping address")
    else:
        # Fallback to using last name as proxy for household
        df_merge_grouped['HOUSEHOLD_GROUP'] = df_merge_grouped['Last_Name'].fillna('Unknown')
        print("Shipping address fields not available, grouping households by last name instead")
    
    # Add household member info
    household_members = {}
    for household_id in df_merge_grouped['HOUSEHOLD_GROUP'].unique():
        members = df_merge_grouped[df_merge_grouped['HOUSEHOLD_GROUP'] == household_id][['CLIENT_UNIQUE_ID']].drop_duplicates()
        member_list = [name for name in members['CLIENT_UNIQUE_ID'] if name.strip()]
        household_members[household_id] = member_list
    
    # Add this info to the dataframe
    df_merge_grouped['HOUSEHOLD_MEMBERS'] = df_merge_grouped['HOUSEHOLD_GROUP'].map(
        lambda x: "\n".join(household_members[x])
    )
    
    # Generate household labels
    df_merge_grouped['line_household_label'] = df_merge_grouped.apply(
        lambda row: ("Household Members:" if len((row['HOUSEHOLD_MEMBERS'] or '').split('\n')) > 1 else ""),
        axis=1
    )
    
    # Generate final text fields
    df_merge_grouped['line_household'] = df_merge_grouped.apply(
        lambda row: (row['HOUSEHOLD_MEMBERS'] if len((row['HOUSEHOLD_MEMBERS'] or '').split('\n')) > 1 else ""),
        axis=1
    )
    
    df_merge_grouped['line_preparedFor'] = df_merge_grouped.apply(
        lambda row: "Prepared for: " + (row['First_Name'] or '') + " " + (row['Last_Name'] or ''), axis=1
    )
    # Add best before date (delivery date + 3 days)
    def add_days(date_str, days=3):
        try:
            date = datetime.strptime(date_str, "%Y-%m-%d")
        except ValueError:
            try:
                date = datetime.fromisoformat(date_str.replace('Z', '+00:00'))
            except ValueError:
                date = datetime.now()
        
        end_date = date + timedelta(days=days)
        return end_date.strftime('%m/%d/%Y')
    
    df_merge_grouped['line_bestBefore'] = "Best before: " + df_merge_grouped['Delivery Date'].apply(lambda x: add_days(x, 3))
    df_merge_grouped['line_deliveryDate'] = df_merge_grouped['Delivery Date'] + " Delivery"
    df_merge_grouped['line_portion'] = df_merge_grouped.portion_list
    df_merge_grouped['line_nutrition'] = df_merge_grouped.NUTRITION
    df_merge_grouped['line_dishes'] = df_merge_grouped.dish_list
    #number of items on the page (count dishes in dish_list, split by '\n\n')
    df_merge_grouped['line_totalItems'] = df_merge_grouped['dish_list'].apply(lambda x: str(len(x.split('\n\n'))) + ' item(s)' if x else '0 item(s)')
    #if has ice pack item, if so, mark true, otherwise false
    df_merge_grouped['line_icePack'] = df_merge_grouped['dish_list'].apply(lambda x: 'Include ice pack' if 'Ice Pack' in str(x) else '')
    df_merge_grouped['line_zone'] = 'Zone ' + df_merge_grouped['ZONE_NUMBER']
    
    
    # Sort by household group to ensure all members of the same household are together
    df_merge_grouped = df_merge_grouped.sort_values(by=['HOUSEHOLD_GROUP', 'First_Name', 'Last_Name', 'page_number'])
    
    # Debug log to show household groupings
    print(f"Total households: {len(df_merge_grouped['HOUSEHOLD_GROUP'].unique())}")
    for household in df_merge_grouped['HOUSEHOLD_GROUP'].unique():
        members = df_merge_grouped[df_merge_grouped['HOUSEHOLD_GROUP'] == household]['CLIENT_UNIQUE_ID'].unique()
        print(f"Household '{household}': {', '.join(members)}")
    
    return df_merge_grouped

def insert_background(slide, img_path):
    left = top = Inches(0)
    pic = slide.shapes.add_picture(img_path, left, top, width=slide.prs.slide_width, height=slide.prs.slide_height)

    # This moves it to the background
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)
    return slide

# Function to copy a slide with images properly
def copy_slide_with_images(source_slide, target_prs):
    """
    Copy a slide including all images and other elements

    Parameters:
    source_slide (Slide): Source slide to copy
    target_prs (Presentation): Target presentation to add the slide to

    Returns:
    Slide: The new slide in the target presentation
    """
    # Add a new slide with the same layout
    target_slide = target_prs.slides.add_slide(source_slide.slide_layout)

    # Keep track of image shapes to process separately
    image_paths = []

    # Copy all shapes
    for shape in source_slide.shapes:
        if hasattr(shape, 'image') and shape.image:
            # For image shapes, extract and save the image temporarily
            image_data = shape.image.blob
            image_path = f"temp_image_{shape.shape_id}.png"
            with open(image_path, "wb") as f:
                f.write(image_data)
            image_paths.append((image_path, shape.left, shape.top, shape.width, shape.height))
        else:
            # For non-image shapes, try to copy the XML element
            try:
                element = shape.element
                newel = copy.deepcopy(element)
                target_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
            except Exception as e:
                print(f"Error copying shape: {e}")

    # Add the images to the new slide
    for image_path, left, top, width, height in image_paths:
        try:
            target_slide.shapes.add_picture(image_path, left, top, width=width, height=height)
        except Exception as e:
            print(f"Error adding image: {e}")
        finally:
            # Clean up temporary file
            if os.path.exists(image_path):
                os.remove(image_path)

    return target_slide

# Function to insert an instruction sheet
def insert_instruction_sheet(prs, instruction_template_path=None):
    """
    Insert an instruction sheet using either a template file or the default template

    Parameters:
    prs (Presentation): PowerPoint presentation to add the slide to
    instruction_template_path (str): Path to the instruction template file

    Returns:
    Presentation: Updated presentation with instruction sheet added
    """
    if instruction_template_path and os.path.exists(instruction_template_path):
        # Check file extension
        file_ext = os.path.splitext(instruction_template_path)[1].lower()
        
        if file_ext == '.pdf':
            # For PDF files, create a slide with a note to refer to the PDF
            blank_slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
            
            # Add title
            if hasattr(blank_slide, 'shapes') and hasattr(blank_slide.shapes, 'title') and blank_slide.shapes.title:
                blank_slide.shapes.title.text = "Cooking Instructions"
            else:
                # Add a title text box if the layout doesn't have a title placeholder
                left = Inches(1)
                top = Inches(0.5)
                width = Inches(8)
                height = Inches(1)
                title_box = blank_slide.shapes.add_textbox(left, top, width, height)
                title_box.text_frame.text = "Cooking Instructions"
            
            # Add instruction text
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(5)
            textbox = blank_slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            p = text_frame.add_paragraph()
            p.text = f"Please refer to the separate PDF file: {os.path.basename(instruction_template_path)}"
            
            # Add another paragraph mentioning the chef
            chef_name = os.path.basename(instruction_template_path).split('_')[-1].split('.')[0]
            p2 = text_frame.add_paragraph()
            p2.text = f"Chef: {chef_name}"
            
        elif file_ext == '.pptx':
            # For PPTX files, copy slides from the template
            instruction_prs = Presentation(instruction_template_path)
            for slide in instruction_prs.slides:
                # Copy each slide from the instruction template
                copy_slide(slide, prs)
        else:
            # For unsupported file types, create a generic slide
            blank_slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
            if hasattr(blank_slide, 'shapes') and hasattr(blank_slide.shapes, 'title') and blank_slide.shapes.title:
                blank_slide.shapes.title.text = "Cooking Instructions"
            
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(5)
            textbox = blank_slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            p = text_frame.add_paragraph()
            p.text = "Standard cooking instructions go here. Please follow these guidelines for all meals."
    else:
        # Create a basic instruction slide if no template is available
        blank_slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
        if hasattr(blank_slide, 'shapes') and hasattr(blank_slide.shapes, 'title') and blank_slide.shapes.title:
            blank_slide.shapes.title.text = "Cooking Instructions"
        
        # Add basic instructions text box
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(5)
        textbox = blank_slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()
        p.text = "Standard cooking instructions go here. Please follow these guidelines for all meals."

    return prs

def generate_ppt(df, template_path):
    prs = Presentation(template_path)
    df.sort_values(by=['HOUSEHOLD_MEMBERS', 'CLIENT_UNIQUE_ID'], inplace=True)
    #df.to_excel("output.xlsx", index=False)
    # Get the first slide as a template
    template_slide = prs.slides[0]
    text_type = template_slide.shapes[1].shape_type  # Identify the text shape type from template

    # Check if there's an instruction slide in the template
    instruction_slide = None
    if len(prs.slides) > 1:
        instruction_slide = prs.slides[1]

    # Save the template slides before we start adding new ones
    template_slides = []
    for i in range(len(prs.slides)):
        template_slides.append(prs.slides[i])

    current_household = None

    # Iterate each row in df_merge
    for index, row in df.iterrows():
        # If we're changing households and there's an instruction slide, add a copy of it
        if current_household is not None and current_household != row['HOUSEHOLD_GROUP'] and instruction_slide is not None:
            # Add the instruction slide between households
            copy_slide_with_images(instruction_slide, prs)
        
        # Update current household
        current_household = row['HOUSEHOLD_GROUP']
        
        # Add one page copied from template
        new_slide = copy_slide_with_images(template_slide, prs)

        # parse the text to the template slide
        for shape in new_slide.shapes:
            if shape.shape_type == text_type and shape.has_text_frame:
                key = shape.name 
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    if text_frame.paragraphs and text_frame.paragraphs[0].runs:
                        text_frame.paragraphs[0].runs[0].text = row[key]

    copy_slide_with_images(instruction_slide, prs)
    # Remove the template slides (first N slides)
    for _ in range(len(template_slides)):
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    return prs

def generate_one_pagers(db,template_path):

    # Process data from Airtable
    processed_data = process_data(db)
    # processed_data.to_excel("output.xlsx", index=False)
    # Generate PowerPoint
    prs = generate_ppt(processed_data, template_path)

    return prs


if __name__ == "__main__":

    # Initialize Airtable client
    db = new_database_access()
    
    # Generate one pagers
    prs = generate_one_pagers(db,template_path = 'template/One_Pager_Template_v2.pptx')
    
    # Save the result
    output_path = f'OnePager_{datetime.now().strftime("%Y%m%d_%H%M%S")}.pptx'
    prs.save(output_path)
    print(f"One pagers generated and saved as {output_path}")