# Standard library imports
import os
import json
import glob
from datetime import datetime, timedelta
# Third-party imports
import barcode
import time
from barcode.writer import ImageWriter
import brother_ql
from brother_ql.conversion import convert
from brother_ql.backends.helpers import send
from dotenv import load_dotenv
import pandas as pd
from PIL import Image
import pytz
from pptx import Presentation
from pptx.util import Inches
from pyairtable.api.table import Table
import streamlit as st
from store_access import new_database_access

# Local imports
import copy
est = pytz.timezone('US/Eastern')
current_date_time = datetime.now(est).strftime("%Y%m%d_%H%M")

# BASE_ID = "appkvckTQ86vhjSC0"  # 041025 
BASE_ID = "appEe646yuQexwHJo" # production
TABLE_ID = "tblVwpvUmsTS2Se51"  # ClientServings
VIEW_ID = "viw5hROs9I9vV0YEq" # Sorted (For Dish Sticker)

# Set the time zone to Eastern Standard Time (EST)
est = pytz.timezone('US/Eastern')

class AirTableError(Exception):
    """Custom exception for AirTable related errors"""
    pass

class PPTGenerationError(Exception):
    """Custom exception for PowerPoint generation errors"""
    pass


def resize_image(image_path, target_width, target_height):
    """
    Resize the image to the specified width and height.

    Args:
        image_path (str): Path to the image file.
        target_width (int): Target width for the resized image.
        target_height (int): Target height for the resized image.
    """
    try:
        with Image.open(image_path) as img:
            img = img.resize((target_width, target_height), Image.Resampling.LANCZOS)
            img.save(image_path)
        print(f"Successfully resized {image_path} to {target_width}x{target_height}")
    except Exception as e:
        print(f"Error resizing {image_path}: {str(e)}")

def ensure_file_saved(file_path, max_retries=5, delay=0.1):
    """
    Ensure a file is completely saved before proceeding.
    
    Args:
        file_path (str): Path to the file to check
        max_retries (int): Maximum number of retry attempts
        delay (float): Delay between retries in seconds
    
    Returns:
        bool: True if file is saved and accessible, False otherwise
    """
    file_path = file_path
    for attempt in range(max_retries):
        try:
            # Check if file exists and is readable
            if os.path.exists(file_path) and os.access(file_path, os.R_OK):
                # Try to open the file to ensure it's not locked
                with open(file_path, 'rb') as f:
                    f.read(1)  # Try to read at least 1 byte
                return True
        except (OSError, IOError):
            pass
        
        time.sleep(delay)
    
    return False
    
def sort_key(filename):
    """
    Extract the numeric part from the filename for sorting.

    Args:
        filename (str): The filename to extract the number from.

    Returns:
        int: The extracted number for sorting.
    """
    last_part = filename.split('-')[-1]
    number = int(last_part.split('.')[0])
    return number

def print_stickers():
    """
    Print the generated dish stickers using Brother QL printer.
    """
    current_file_path = os.path.abspath(__file__)
    
    # Set up printer configuration
    model = 'QL-810W'
    
    # Configure source folder and printer settings
    source_folder = "/Users/clairegoldwitz/Desktop/image_test/test.jpg"
    backend = 'pyusb'
    identifier = 'usb://0x04f9:0x209c'
    
    # Alternative configuration for network printing
    # source_folder = os.path.join(os.path.dirname(current_file_path), 'assets', 'images')
    # backend = 'network'
    # identifier = 'tcp://192.168.4.30'
    
    # Load all JPG and PNG files in order
    try:
        label_files = glob.glob(f'{source_folder}/*.jpg') + glob.glob(f'{source_folder}/*.png')
        label_files.sort(key=sort_key, reverse=True)
        print(f"Found {len(label_files)} image files")
    except Exception as e:
        print(f"Error loading image files: {str(e)}")
        return
    
    if not label_files:
        print("No image files found in the directory.")
        return
    
    for i, label_file in enumerate(label_files):
        # print(f'LABEL FILE: {label_file} | {i+1} of {len(label_files)}')
        # Create the label maker object
        try:
            qlr = brother_ql.BrotherQLRaster(model)
            qlr.exception_on_warning = True
            print("Successfully created label maker object")
        except Exception as e:
            print(f"Error creating label maker object: {str(e)}")
            return
        
        img = Image.open(label_file)
        # img = img.resize((3800, 2512))
        
        try:
            instructions = convert(
                qlr=qlr,
                images=[img],
                label='62',
                cut=True,
                rotate='90'
            )
            print(f"Successfully converted label {i+1}")
        except Exception as e:
            print(f"Error converting label {i+1}: {str(e)}")
            return
        
        try:
            send(instructions=instructions, printer_identifier=identifier, backend_identifier=backend)
            print(f"Successfully sent label {i+1} to printer")
        except Exception as e:
            print(f"Error sending label {i+1} to printer: {e}")
            return
            
        # Delete previous file:
        try:
            if i > 0 and os.path.exists(label_files[i - 1]):
                print(f"Deleting previous file: {label_files[i - 1]}")
                os.remove(label_files[i - 1])
        except Exception as e:
            print(f"Error deleting previous file: {e}. Check permissions.")
            
    try:
        if os.path.exists(label_files[-1]):
            print(f"Deleting previous file: {label_files[-1]}")
            os.remove(label_files[-1])
    except Exception as e:
        print(f"Error deleting previous file: {e}. Check permissions.")
            
    print("\nPrinting process completed.")

# Process Data 
def generate_sticker_df(df):
    df_dish = df.fillna('N/A')
    
    # Generate Requested Columns
    # barcode ID: got a float - change to string
    df_dish["#"] = df_dish["#"].apply(lambda x: str(int(x)))

    # all linked item needs to be extracted from a list
    df_dish['client_name'] = df_dish['Customer Name'].apply(lambda x: x[0] if isinstance(x, list) else x)

    # breakfast, lunch, etc.
    df_dish['meal'] = df_dish['Meal Portion (from Linked OrderItem)'].apply(lambda x: x[0] if isinstance(x, list) else x)

    # zone
    df_dish['zone'] = df_dish['Delivery Zone (from Linked OrderItem)'].apply(lambda x: 'Zone ' + str(x[0]) if isinstance(x, list) else 'Zone ' + str(x))

    def add_days(date_str, days = 3):  # by default: add 3 days based on EST
        date = datetime.now(est).strptime(date_str, "%Y-%m-%d")
        end_date = date + timedelta(days)
        return end_date.strftime('%m/%d/%Y')
    
    df_dish['Best Before'] = df_dish['Delivery Date'].apply(lambda x: add_days(x[0] if isinstance(x, list) else x, 3))
    df_dish['best_before'] = "Best before: " + df_dish['Best Before']

    meal_info = df_dish['Meal Sticker (from Linked OrderItem)'].apply(lambda x: x[0] if isinstance(x, list) else x)
    df_dish[['bowl_name', 'ingredients']] = meal_info.str.split(': ', n=1, expand=True)
    df_dish['parts'] = df_dish['# of Parts'].apply(lambda x: x[0] if isinstance(x, list) else x)

    # duplicate according to '# portions'
    # df_dish = df_dish.loc[df_dish.index.repeat(df_dish['# portions'])]
    # df_dish = df_dish.reset_index(drop=True)
    
    return df_dish

"""### Draw Slides"""

def insert_background(slide, img_path, prs):
    left = top = Inches(0)
    pic = slide.shapes.add_picture(img_path, left, top, 
                                width=prs.slide_width, height=prs.slide_height)
    # This moves it to the background
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)
    return slide

def copy_slide(template, prs):
    # Create a new slide with the same layout as the source slide
    new_slide = prs.slides.add_slide(template.slide_layout)

    # Copy the content from the source slide to the new slide
    for shape in template.shapes:
        new_shape = copy.deepcopy(shape)
        new_slide.shapes._spTree.insert_element_before(new_shape._element, 'p:extLst')

    return new_slide

def generate_dish_stickers_barcode(db):
    client_serving_df = read_client_serving(db)
    # print(df.shape)
    # print(df.columns)
    df = generate_sticker_df(client_serving_df)
    # Load the presentation
    prs = Presentation('template/Dish_Sticker_Template_Barcode.pptx')

    # Get the first slide
    slide = prs.slides[0]
    # for i in slide.shapes:
    #     print('%d %s' % (i.shape_id, i.name))
    #     print(i.text_frame.text)

    # https://python-barcode.readthedocs.io/en/stable/getting-started.html#usage
    ITF = barcode.get_barcode_class('ITF')
    BARCODE_HEIGHT = 0.75
    margin = Inches(0.05)

    
    if 'Dish ID (from Linked OrderItem)' in df.columns and 'Position Id' in df.columns:
        df.sort_values(by=['Dish ID (from Linked OrderItem)', 'Position Id'], ascending=[True, True], inplace=True)

    # Find the maximum number of parts across all rows
    max_parts = df['parts'].max()

    # Iterate by part number first (all Part 1s, then all Part 2s, etc.)
    for part_num in range(max_parts):
        # iterate each row in df_dish
        for _, row in df.iterrows():
            parts = row['parts']
            # Skip if this row doesn't have this many parts
            if part_num >= parts:
                continue

            # add one page copied from template
            new_slide = copy_slide(slide, prs)

            for shape in new_slide.shapes:
                if shape.has_text_frame:
                    key = shape.name
                    text_frame = shape.text_frame
                    if text_frame.paragraphs and text_frame.paragraphs[0].runs:
                        value = row.get(key, 'N/A')
                        if key == 'parts':
                            if parts > 1:
                                value = f"PART {part_num+1}/{parts}"
                            else:
                                value = ''
                            #print(f"DEBUG: key={key}, parts={parts}, value={value}")
                        text_frame.paragraphs[0].runs[0].text = str(value) if value is not None else 'N/A'

            # ITF barcode generator, with transparent background
            fileName = f'id_barcode_{row["#"]}'
            itf = ITF(row["#"], writer = ImageWriter(mode = "RGBA"))
            itf.save(fileName, dict(quiet_zone=3, background = (255, 255, 255, 0),
                                        font_size = 5, text_distance = 2.5,
                                        module_width = 0.3))

            # Ensure file is saved before proceeding
            if not ensure_file_saved(fileName + '.png'):
                raise PPTGenerationError(f"Failed to save barcode file: {fileName}")

            id_barcode = Image.open(fileName + '.png')
            barcode_size = ((Inches(id_barcode.size[0] /id_barcode.size[1] * BARCODE_HEIGHT)), Inches(BARCODE_HEIGHT))

            new_slide.shapes.add_picture(fileName + '.png', prs.slide_width - margin - barcode_size[0],
                                            prs.slide_height - margin - barcode_size[1],
                                            width = barcode_size[0], height = barcode_size[1])
            # delete the file
            os.remove(fileName + '.png')
    return prs

def read_weekly_products(db):
    fields_to_return = ['Internal Dish ID', '# of Parts']
    try:
        records = db.get_weekly_products()
        if not records:
            raise AirTableError("No records found in the specified view.")
        
        # Convert records to DataFrame
        df = pd.DataFrame(records)
        
        # Extract fields from the records
        if 'fields' in df.columns:
            fields_df = pd.json_normalize(df['fields'])
            # Combine the records with their fields
            df_flat = pd.concat([df.drop('fields', axis=1), fields_df], axis=1)
        else:
            df_flat = df

        df_flat = df_flat[fields_to_return]

        return df_flat
    except Exception as e:
        raise AirTableError(f"Error fetching data from Airtable: {str(e)}")

def read_client_serving(db):
    fields_to_return = ['#', 'Customer Name', 'Meal Sticker (from Linked OrderItem)', 'Meal Portion (from Linked OrderItem)', 'Delivery Date', 'Position Id', 'Dish ID (from Linked OrderItem)','Delivery Zone (from Linked OrderItem)','# of Parts']
    try:
        # Initialize table
        records = db.get_clientservings_data(view=VIEW_ID)
        
        if not records:
            raise AirTableError("No records found in the specified view.")
        
        # Convert records to DataFrame
        df = pd.DataFrame(records)
        
        # Extract fields from the records
        if 'fields' in df.columns:
            fields_df = pd.json_normalize(df['fields'])
            # Combine the records with their fields
            df_flat = pd.concat([df.drop('fields', axis=1), fields_df], axis=1)
        else:
            df_flat = df
        
        if 'Dish ID (from Linked OrderItem)' in df_flat.columns and 'Position Id' in df_flat.columns:
            df_flat['Dish ID (from Linked OrderItem)'] = df_flat['Dish ID (from Linked OrderItem)'].apply(
                lambda x: int(x[0]) if isinstance(x, list) and len(x) > 0 and str(x[0]).isdigit() 
                else (int(x) if isinstance(x, (int, float)) else x)
            )
        df_flat = df_flat[fields_to_return]
            
        return df_flat
        
    except Exception as e:
        raise AirTableError(f"Error fetching data from Airtable: {str(e)}")

if __name__ == "__main__":
    db = new_database_access()

    prs = generate_dish_stickers_barcode(db)
    prs.save(f'{current_date_time}_dish_sticker_barcode.pptx')

    #print_stickers()