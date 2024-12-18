# -*- coding: utf-8 -*-
import streamlit as st
import pandas as pd
from pptx import Presentation
import copy
from datetime import datetime
from math import ceil
from pptx.util import Pt

def process_shipping_data(uploaded_shipping_file):
    df_shipping = pd.read_csv(uploaded_shipping_file,encoding_errors='ignore')
    df_shipping = df_shipping.astype(str)
    df_shipping['Shipping Phone'] = df_shipping['Shipping Phone'].str.replace(r'[()-]', '', regex=True)
    df_shipping.fillna('', inplace=True)
    df_shipping = df_shipping.replace('nan', '')
    df_shipping = df_shipping.replace('-', '')
    df_shipping['Shipping Province'] = df_shipping['Shipping Province'].str.upper()
    df_shipping['merge_key_shippingname'] = df_shipping['Shipping Name'].str.lower()
    return df_shipping

def fetch_client_servings(sheet_id, sheet_name):
    url = f"https://docs.google.com/spreadsheets/d/{sheet_id}/gviz/tq?tqx=out:csv&sheet={sheet_name}"
    df = pd.read_csv(url)
    df = df[['Client', 'Meal', '# portions']].dropna().copy()
    df['Client'] = df['Client'].str.replace(r'\s*\((.*?)\)', r', \1', regex=True)
    df.loc[df['Meal'].isin(['Breakfast', 'Snack']), '# portions'] *= 0.5
    df = df.groupby('Client')['# portions'].sum().reset_index(name='total_portion_per_pp')
    return df

def fetch_package_recipient(sheet_id, sheet_name):
    url = f"https://docs.google.com/spreadsheets/d/{sheet_id}/gviz/tq?tqx=out:csv&sheet={sheet_name}"
    df = pd.read_csv(url)
    df = df[['Name', 'Other Member of Household']].dropna().copy()
    df.drop_duplicates(inplace=True)
    return df

def add_recipient_clientservings(all_client_servings, client_df):
    df=all_client_servings.merge(client_df,left_on='Client',right_on='Name',how='left')
    df.rename(columns={'Other Member of Household': 'Package recipient'},inplace=True)
    df['Package recipient'] = df['Package recipient'].fillna(df['Client'])
    df.drop(columns=['Name'], inplace=True)
    df['merge_key_PcgRcpt'] = df['Package recipient'].str.lower()
    df['merge_key_Client'] = df['Client'].str.lower()
    return df


def match_orders_to_shipping_data(client_servings_df, shipping_df):
    portion_per_sticker=6
    match_df = client_servings_df.merge(shipping_df, left_on='merge_key_PcgRcpt', right_on='merge_key_shippingname', how='left')
    match_df = match_df.dropna(subset=['merge_key_shippingname'])
    match_df['Shipping Phone'] = match_df['Shipping Phone'].astype(str).str[:3] + '-' + match_df['Shipping Phone'].astype(str).str[3:6] + '-' + match_df['Shipping Phone'].astype(str).str[6:10]
    match_df.drop(['merge_key_shippingname', 'merge_key_PcgRcpt', 'merge_key_Client'], axis=1, inplace=True)
    final_df = match_df.groupby(['Shipping Name', 'Shipping Address1', 'Shipping Address2', 'Shipping City', 'Shipping Province', 'Shipping Zip', 'Shipping Phone'])['total_portion_per_pp'].sum().reset_index(name='total_portion_per_address')
    final_df['#_of_stickers'] = final_df['total_portion_per_address'].apply(lambda x: ceil(x / portion_per_sticker) * 2)
    return final_df

# Function to copy a slide from a template
def copy_slide(template, target_prs):
    new_slide = target_prs.slides.add_slide(template.slide_layout)
    for shape in new_slide.shapes:
        sp = shape._element
        sp.getparent().remove(sp)
    for shape in template.shapes:
        if not shape.has_text_frame:
            continue
        new_shape = copy.deepcopy(shape)
        new_slide.shapes._spTree.insert_element_before(new_shape._element, 'p:extLst')
    return new_slide

def generate_ppt(final_match_result_with_portion_df, prs):
    
    last_slide_index = 0
    template_slide = prs.slides[0]
    for index, row in final_match_result_with_portion_df.iterrows():
        slide = copy_slide(template_slide, prs)
        last_slide_index += 1
        sticker_needed = row['#_of_stickers']
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for i, paragraph in enumerate(text_frame.paragraphs):
                    if 'Shipping Name' in paragraph.text:
                        paragraph.text = row['Shipping Name']
                        paragraph.font.size = Pt(28)
                        paragraph.font.name = "Calibri"
                    elif "Address" in paragraph.text:
                        if row['Shipping Address2'] != '':
                            paragraph.text = f"{row['Shipping Address1']}, {row['Shipping Address2']}"
                            paragraph.font.size = Pt(24)
                            paragraph.font.name = "Calibri"
                        else:
                            paragraph.text = f"{row['Shipping Address1']}"
                            paragraph.font.size = Pt(24)
                            paragraph.font.name = "Calibri"
                    elif "City" in paragraph.text:
                        paragraph.text = f"{row['Shipping City']}, {row['Shipping Province']} {row['Shipping Zip']}"
                        paragraph.font.size = Pt(24)
                        paragraph.font.name = "Calibri"
                    elif 'Shipping Phone' in paragraph.text:
                        paragraph.text = str(row['Shipping Phone'])
                        paragraph.font.size = Pt(24)
                        paragraph.font.name = "Calibri"
        if sticker_needed > 1:
            count = 1
            while count < sticker_needed:
                copy_slide(prs.slides[last_slide_index], prs)
                last_slide_index += 1
                count += 1
    return prs

if __name__ == "__main__":
    sheet_id = "1rorOBlH_K9qH4L39KehvI_rYGHo7agNVdCDisWydEj8"
    LA_sheet_name = "ClientServings-LA"
    NY_sheet_name = "ClientServings"
    client_sheet_name = "Clients"

    LA_clientservings_df = fetch_client_servings(sheet_id, LA_sheet_name)
    
    NY_clientservings_df = fetch_client_servings(sheet_id, NY_sheet_name)
    print(NY_clientservings_df)