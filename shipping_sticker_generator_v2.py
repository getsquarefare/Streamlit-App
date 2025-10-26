from pptx import Presentation
from pptx.util import Pt
from datetime import datetime
from math import ceil
from io import BytesIO
import logging
import copy
from exceptions import AirTableError
from store_access import new_database_access

# Set up logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

VIEW = "viwdOQZFrb6gEzkU9"  # View for open orders

class PPTGenerationError(Exception):
    """Custom exception for PowerPoint generation errors"""
    pass


def process_order_data(db):
    """Process order data from Airtable and group by shipping address"""
    try:
        logger.info("Starting order data processing")

        # Get data from Airtable
        orders = db.get_all_open_orders(view=VIEW)

        if not orders:
            raise AirTableError("No open orders found")

        logger.info(f"Found {len(orders)} orders")

        # Process each order
        shipping_data = []
        portion_per_sticker = 6.8

        for order in orders:
            fields = order['fields']

            # Skip orders with missing required fields
            required_fields = ['Shipping Name', 'Shipping Address 1', 'Quantity']
            if not all(fields.get(field) for field in required_fields):
                logger.warning(f"Skipping order with missing fields: {order.get('id', 'unknown')}")
                continue

            # Calculate quantity based on meal portion
            meal_portion = fields.get('Meal Portion', '')
            quantity = fields.get('Quantity', 0)

            if meal_portion == 'Breakfast':
                adjusted_quantity = quantity * 0.8
            elif meal_portion == 'Snack':
                adjusted_quantity = quantity * 0.1
            else:
                adjusted_quantity = quantity

            # Format phone number
            phone = fields.get('Shipping Phone', '')
            if phone:
                # Remove non-digits
                phone_digits = ''.join(filter(str.isdigit, str(phone)))
                if len(phone_digits) == 10:
                    phone = f"{phone_digits[:3]}-{phone_digits[3:6]}-{phone_digits[6:]}"

            # Extract zone number from list
            zone = fields.get('Zone Number (from Delivery Zone)', 'N/A')
            if isinstance(zone, list) and len(zone) > 0:
                zone = str(zone[0])

            # Create shipping record
            shipping_record = {
                'Shipping Name': fields.get('Shipping Name', ''),
                'Shipping Address 1': fields.get('Shipping Address 1', ''),
                'Shipping Address 2': fields.get('Shipping Address 2', ''),
                'Shipping City': fields.get('Shipping City', ''),
                'Shipping Province': fields.get('Shipping Province', '').upper(),
                'Shipping Postal Code': fields.get('Shipping Postal Code', ''),
                'Shipping Phone': phone,
                'Zone Number': zone,
                'Quantity': adjusted_quantity
            }

            shipping_data.append(shipping_record)

        # Group by shipping address
        grouped_shipping = {}
        for record in shipping_data:
            # Create unique key for grouping
            key = (
                record['Shipping Name'],
                record['Shipping Address 1'],
                record['Shipping Address 2'],
                record['Shipping City'],
                record['Shipping Province'],
                record['Shipping Postal Code'],
                record['Shipping Phone'],
                record['Zone Number']
            )

            if key not in grouped_shipping:
                grouped_shipping[key] = {
                    'Shipping Name': record['Shipping Name'],
                    'Shipping Address 1': record['Shipping Address 1'],
                    'Shipping Address 2': record['Shipping Address 2'],
                    'Shipping City': record['Shipping City'],
                    'Shipping Province': record['Shipping Province'],
                    'Shipping Postal Code': record['Shipping Postal Code'],
                    'Shipping Phone': record['Shipping Phone'],
                    'Zone Number': record['Zone Number'],
                    'Total Quantity': 0
                }

            grouped_shipping[key]['Total Quantity'] += record['Quantity']

        # Calculate number of stickers needed for each address
        shipping_list = []
        for shipping_info in grouped_shipping.values():
            stickers_needed = ceil(shipping_info['Total Quantity'] / portion_per_sticker) * 2
            shipping_info['Stickers Needed'] = stickers_needed
            shipping_list.append(shipping_info)

        logger.info(f"Processed {len(shipping_list)} unique shipping addresses")
        return shipping_list

    except Exception as e:
        logger.error(f"Error processing order data: {str(e)}")
        raise AirTableError(f"Error processing order data: {str(e)}")

def copy_slide(template_slide, target_prs):
    """Copy a slide from template to target presentation"""
    try:
        new_slide = target_prs.slides.add_slide(template_slide.slide_layout)

        # Remove default shapes
        for shape in new_slide.shapes:
            sp = shape._element
            sp.getparent().remove(sp)

        # Copy shapes from template
        for shape in template_slide.shapes:
            if not shape.has_text_frame:
                continue
            new_shape = copy.deepcopy(shape)
            new_slide.shapes._spTree.insert_element_before(new_shape._element, 'p:extLst')

        return new_slide

    except Exception as e:
        logger.error(f"Error copying slide: {str(e)}")
        raise PPTGenerationError(f"Failed to copy slide template: {str(e)}")

def populate_sticker(slide, shipping_info):
    """Populate a single sticker slide with shipping information"""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            text = paragraph.text

            if 'Shipping Name' in text:
                paragraph.text = shipping_info['Shipping Name']
                paragraph.font.size = Pt(28)
                paragraph.font.name = "Lato"

            elif "Address" in text:
                if shipping_info['Shipping Address 2']:
                    paragraph.text = f"{shipping_info['Shipping Address 1']}, {shipping_info['Shipping Address 2']}"
                else:
                    paragraph.text = shipping_info['Shipping Address 1']
                paragraph.font.size = Pt(24)
                paragraph.font.name = "Lato"

            elif "City" in text:
                paragraph.text = f"{shipping_info['Shipping City']}, {shipping_info['Shipping Province']} {shipping_info['Shipping Postal Code']}"
                paragraph.font.size = Pt(24)
                paragraph.font.name = "Lato"

            elif 'Shipping Phone' in text:
                paragraph.text = str(shipping_info['Shipping Phone'])
                paragraph.font.size = Pt(24)
                paragraph.font.name = "Lato"

            elif 'ZONE' in text:
                paragraph.text = f"ZONE {shipping_info['Zone Number']}"
                paragraph.font.size = Pt(28)
                paragraph.font.name = "Lato"

def create_shipping_stickers_ppt(shipping_list, template_path='template/Shipping_Sticker_Template.pptx'):
    """Create PowerPoint presentation with shipping stickers"""
    try:
        logger.info("Starting PowerPoint generation")

        if not shipping_list:
            raise PPTGenerationError("No shipping data available to generate stickers")

        # Load template
        prs = Presentation(template_path)

        if len(prs.slides) == 0:
            raise PPTGenerationError("Template presentation has no slides")

        template_slide = prs.slides[0]
        total_stickers = 0

        logger.info(f"Generating stickers for {len(shipping_list)} shipping addresses")

        for shipping_info in shipping_list:
            stickers_needed = shipping_info['Stickers Needed']

            # Create stickers for this address
            for _ in range(stickers_needed):
                slide = copy_slide(template_slide, prs)
                populate_sticker(slide, shipping_info)
                total_stickers += 1

        # Remove the template slide
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

        logger.info(f"Successfully generated {total_stickers} stickers for {len(shipping_list)} addresses")

        # Save to BytesIO
        output = BytesIO()
        prs.save(output)
        output.seek(0)

        return output

    except Exception as e:
        logger.error(f"Error creating PowerPoint: {str(e)}")
        raise PPTGenerationError(f"Error creating shipping stickers: {str(e)}")

def generate_shipping_stickers(db):
    """Main method to generate shipping stickers"""
    template_path = 'template/Shipping_Sticker_Template.pptx'
    try:
        logger.info("Starting shipping sticker generation")

        # Process order data
        shipping_list = process_order_data(db)

        if not shipping_list:
            raise AirTableError("No shipping addresses to process")

        logger.info(f"Processing {len(shipping_list)} shipping addresses")

        # Create PowerPoint file
        ppt_file = create_shipping_stickers_ppt(shipping_list, template_path)

        return ppt_file

    except Exception as e:
        logger.error(f"Error generating shipping stickers: {str(e)}")
        raise AirTableError(f"Error generating shipping stickers: {str(e)}")


if __name__ == "__main__":
    try:
        db = new_database_access()
        ppt_file = generate_shipping_stickers(db)

        # Write the BytesIO content to a file
        filename = f'shipping_stickers_{datetime.now().strftime("%Y%m%d_%H%M%S")}.pptx'
        with open(filename, 'wb') as f:
            f.write(ppt_file.getvalue())

        logger.info(f"Successfully generated shipping stickers: {filename}")

    except AirTableError as e:
        logger.critical(f"Application error: {str(e)}")
    except Exception as e:
        logger.critical(f"Unexpected error: {str(e)}")
