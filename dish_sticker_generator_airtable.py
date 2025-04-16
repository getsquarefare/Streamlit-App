import streamlit as st
import pandas as pd
from pptx import Presentation
import copy
from datetime import datetime, timedelta
from pptx.util import Inches, Pt
from pyairtable.api.table import Table
from pyairtable.formulas import match
import os
from dotenv import load_dotenv
from functools import cache
import traceback
from requests.exceptions import RequestException
from pptx.exc import PackageNotFoundError

class AirTableError(Exception):
    """Custom exception for AirTable related errors"""
    pass

class DataProcessingError(Exception):
    """Custom exception for data processing errors"""
    pass

class PPTGenerationError(Exception):
    """Custom exception for PowerPoint generation errors"""
    pass

class AirTable():
    def __init__(self, ex_api_key=None):
        try:
            # Load environment variables from the .env file
            load_dotenv()
            
            # Get the API key from environment variables or the passed argument
            self.api_key = ex_api_key or st.secrets.get("AIRTABLE_API_KEY")
            if not self.api_key:
                raise AirTableError("Airtable API key not found. Please check your .env file or streamlit secrets.")
                
            self.base_id = "appEe646yuQexwHJo"
            
            # Initialize tables for dish stickers
            self.dish_orders_table = Table(self.api_key, self.base_id, 'tblxT3Pg9Qh0BVZhM')
            
            # Define field mappings for dish stickers
            self.fields = {
                'CLIENT': 'fldWYJStYSpX72pG3',
                'DISH_STICKER': 'fldYZYDRjScz6ig5a',
                'DELIVERY_DAY': 'flddRJziNBdEtrpmG',
                'MEAL_TYPE': 'fldCsBzoy9rxKlWmN',
                'QUANTITY': 'fldvkwFMlBOW5um2y',
                'POSITION_ID':'fldu8pc7PsGFMRRqv',
                'DISH_NAME':'fldegrwFVynxhDavT',
                'DISH_ID':'fldLOvWuvg6X9Odvw'
            }
        except Exception as e:
            raise AirTableError(f"Failed to initialize AirTable connection: {str(e)}")
    
    def get_all_dish_orders(self): 
        try:
            data = self.dish_orders_table.all(fields=self.fields.values(), view='viwxYRtkq19KcXcZX')
            
            if not data:
                raise AirTableError("No data retrieved from Airtable. Please check if the view has records.")
                
            df = pd.DataFrame([record.get('fields', {}) for record in data])
            
            if df.empty:
                raise AirTableError("Retrieved data is empty. Please check your Airtable data.")
            
            # Rename columns to match the expected format
            column_mapping = {v.replace('fld', ''): k for k, v in self.fields.items()}
            df.rename(columns=column_mapping, inplace=True)

            # Convert list values to strings
            df = df.applymap(lambda x: x[0] if isinstance(x, list) and len(x) == 1 else ', '.join(map(str, x)) if isinstance(x, list) else x)

            # Verify required columns
            required_columns = ['Customer Name', 'Order Type', 'Delivery Date', 'Meal Sticker', 'Quantity']
            missing_columns = [col for col in required_columns if col not in df.columns]
            
            if missing_columns:
                raise AirTableError(f"Missing required columns or no data in {', '.join(missing_columns)}")

            return df
            
        except RequestException as e:
            raise AirTableError(f"Network error while connecting to Airtable: {str(e)}")
        except Exception as e:
            raise AirTableError(f"Failed to retrieve dish orders: {str(e)}")

    def process_data(self, df):
        try:
            # Clean column names
            df.columns = [col.strip() for col in df.columns]
            
            # Validate data
            if 'Delivery Date' not in df.columns:
                raise DataProcessingError("Missing 'Delivery Date' column in data")
            if 'Quantity' not in df.columns:
                raise DataProcessingError("Missing 'Quantity' column in data")
                
            # Check for missing values in critical fields
            missing_delivery_dates = df['Delivery Date'].isnull().sum()
            if missing_delivery_dates > 0:
                raise DataProcessingError(f"Found {missing_delivery_dates} records with missing delivery dates")
                
            # Create required fields for dish stickers
            df['line_1'] = df.apply(lambda row: f"Prepared for: {row.get('Customer Name') or 'N/A'} - {row.get('Order Type') or 'N/A'}", axis=1)
            
            # Add best before date (delivery date + 3 days)
            def add_days(date_str, days=3):
                try:
                    date = datetime.strptime(date_str, "%Y-%m-%d")
                    end_date = date + timedelta(days=days)
                    return end_date.strftime('%m-%d-%Y')
                except ValueError:
                    raise DataProcessingError(f"Invalid date format: '{date_str}'. Expected YYYY-MM-DD")
            
            # Safely process dates
            df['Best Before'] = df['Delivery Date'].apply(
                lambda x: add_days(x, 3) if isinstance(x, str) and x else "Date missing"
            )
            df['line_2'] = "Best before: " + df['Best Before']
            
            # Set dish name
            df['line_3'] = df.apply(
                lambda row: row.get('Meal Sticker') or "Dish name not available", 
                axis=1
            )

            # Convert lines to strings
            df['line_1'] = df['line_1'].astype(str)
            df['line_2'] = df['line_2'].astype(str)
            df['line_3'] = df['line_3'].astype(str)
            
            # Safely convert quantity to numeric
            df['Quantity'] = pd.to_numeric(df['Quantity'], errors='coerce').fillna(1).astype(int)
            
            # Duplicate rows based on portion count
            result_df = df.loc[df.index.repeat(df['Quantity'])].reset_index(drop=True)
            if result_df.empty:
                raise DataProcessingError("No data after processing. Check if quantities are valid.")
                
            return result_df
            
        except ValueError as e:
            raise DataProcessingError(f"Error in data values: {str(e)}")
        except Exception as e:
            raise DataProcessingError(f"Failed to process data: {str(e)}")

def insert_background(slide, img_path, prs):
    try:
        if not os.path.exists(img_path):
            raise PPTGenerationError(f"Background image not found: {img_path}")
            
        left = top = Inches(0)
        width = prs.slide_width
        height = prs.slide_height
        pic = slide.shapes.add_picture(img_path, left, top, width=width, height=height)

        # This moves it to the background
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)
        return slide
    except Exception as e:
        raise PPTGenerationError(f"Failed to insert background: {str(e)}")

def copy_slide(template, target_prs):
    try:
        new_slide = target_prs.slides.add_slide(template.slide_layout)
        for shape in new_slide.shapes:
            sp = shape._element
            sp.getparent().remove(sp)
        for shape in template.shapes:
            if not shape.has_text_frame:
                continue
            new_shape = copy.deepcopy(shape)
            new_slide.shapes._spTree.insert_element_before(new_shape._element, 'p:extLst')
        return new_slide
    except Exception as e:
        raise PPTGenerationError(f"Failed to copy slide template: {str(e)}")

def generate_ppt(df, prs, background_path):
    try:
        if df.empty:
            raise PPTGenerationError("No data available to generate stickers")
            
        # Verify required columns
        required_columns = ['line_1', 'line_2', 'line_3', 'Dish ID', 'Position Id (from ClientServings)']
        missing_columns = [col for col in required_columns if col not in df.columns]
        if missing_columns:
            # Handle case where some columns are missing but we can continue
            st.warning(f"Missing or no data in columns: {', '.join(missing_columns)}. Will ignore.")
            for col in missing_columns:
                if col in ['Dish ID', 'Position Id']:
                    df[col] = 'Unknown'  # Default values for sorting
                else:
                    df[col] = 'Not available'  # Default values for text
        
        # Safely sort values
        if 'Dish ID' in df.columns and 'Position Id' in df.columns:
            df.sort_values(by=['Dish ID', 'Position Id'], ascending=[True, True], inplace=True)
            
        # Get the first slide as a template
        if len(prs.slides) == 0:
            raise PPTGenerationError("PowerPoint template has no slides")
            
        template_slide = prs.slides[0]
        
        processed_count = 0
        total_rows = len(df)
        
        # iterate each row in df_dish
        for index, row in df.iterrows():
            # add one page copied from template
            new_slide = copy_slide(template_slide, prs)

            # change text for lines 1-3
            for i, shape in enumerate(new_slide.shapes):
                if shape.has_text_frame and i < 3:  # Only process the first 3 shapes with text
                    # Get the text frame
                    text_frame = shape.text_frame

                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            line_key = f'line_{i+1}'
                            run.text = row.get(line_key, "Text not available")
                        break

            # add background if provided
            if background_path:
                if os.path.exists(background_path):
                    new_slide = insert_background(new_slide, background_path, prs)
                else:
                    st.warning(f"Background image not found: {background_path}")
                    
            processed_count += 1
        
        st.success(f"Successfully generated {processed_count} stickers")
        return prs
    except Exception as e:
        raise PPTGenerationError(f"Error generating PowerPoint: {str(e)}")

def new_database_access():
    try:
        return AirTable()
    except AirTableError as e:
        st.error(f"Database connection error: {str(e)}")
        st.stop()

def generate_dish_stickers(prs, background_path=None):
    try:
        ac = new_database_access()
        data = ac.get_all_dish_orders()
        processed_data = ac.process_data(data)
        
        st.success(f"Processed {len(processed_data)} stickers to generate")
        prs = generate_ppt(processed_data, prs, background_path)
        
        return prs
    except AirTableError as e:
        st.error(f"Airtable Error: {str(e)}")
        st.info("Please check your API key and Airtable connection settings")
        return None
    except DataProcessingError as e:
        st.error(f"Data Processing Error: {str(e)}")
        st.info("Please verify the data in your Airtable and ensure all required fields are present")
        return None
    except PPTGenerationError as e:
        st.error(f"PowerPoint Generation Error: {str(e)}")
        st.info("Please check your PowerPoint template file and try again")
        return None
    except Exception as e:
        st.error(f"Unexpected Error: {str(e)}")
        st.info("Error details:")
        st.code(traceback.format_exc())
        return None
