#!/usr/bin/env python3
"""
Read bag_barcode_mapping.json and generate a PowerPoint preview using the
same template + layout as production bag/shipping stickers.

Does NOT read Open Orders or write to Airtable — JSON only.

Examples:
  cd Streamlit-App && .venv/bin/python scripts/preview_bag_stickers.py
  .venv/bin/python scripts/preview_bag_stickers.py --json bag_barcode_mapping.json
  .venv/bin/python scripts/preview_bag_stickers.py --bag BAG-20260518-Z3-0001 -o /tmp/one_bag.pptx
  .venv/bin/python scripts/preview_bag_stickers.py --limit 3
"""

from __future__ import annotations

import argparse
import json
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation

BASE_DIR = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(BASE_DIR))

from src.stickers.generate_bag_stickers_barcode import (
    DEFAULT_BAG_TEMPLATE,
    _uses_one_pager_template,
    copy_slide_with_images,
    populate_bag_slide,
)
from src.stickers.shipping_sticker_generator_v3 import add_code128_barcode, copy_slide, format_phone

DEFAULT_JSON = BASE_DIR / "bag_barcode_mapping.json"


def find_json_candidates(root: Path) -> list[Path]:
    names = ("bag_barcode_mapping.json", "bag_barcode_mapping*.json")
    found: list[Path] = []
    for pattern in names:
        found.extend(root.glob(pattern))
        found.extend(root.parent.glob(pattern))
    unique: list[Path] = []
    seen: set[Path] = set()
    for p in sorted(found, key=lambda x: x.stat().st_mtime, reverse=True):
        rp = p.resolve()
        if rp not in seen and p.is_file():
            seen.add(rp)
            unique.append(p)
    return unique


def load_mapping(path: Path, inject_mock_phone: bool = False) -> list[dict]:
    with path.open(encoding="utf-8") as f:
        data = json.load(f)
    if not isinstance(data, list):
        raise ValueError(f"{path} must be a JSON array of bag objects")
    if inject_mock_phone:
        for i, bag in enumerate(data):
            if not str(bag.get("phone", "")).strip():
                bag["phone"] = f"212-555-{1000 + i:04d}"
    return data


def json_bag_to_slide_row(bag: dict) -> dict:
    """Map saved JSON entry → row dict used by populate_bag_slide / populate_sticker."""
    dishes = bag.get("dishes") or []
    dish_lines: list[str] = []
    for d in dishes:
        if isinstance(d, dict):
            line = (
                d.get("displayText")
                or " - ".join(
                    x
                    for x in [
                        d.get("customerName"),
                        d.get("mealPortion"),
                        d.get("mealSticker"),
                    ]
                    if x
                )
            )
            dish_lines.append(str(line).strip())
        else:
            dish_lines.append(str(d).strip())
    dish_list = "\n\n".join(x for x in dish_lines if x)

    members = bag.get("householdMembers") or []
    if isinstance(members, str):
        members = [members]
    household = "\n".join(str(m).strip() for m in members if str(m).strip())

    zone = str(bag.get("zone", "")).replace("Zone", "").replace("ZONE", "").strip()
    try:
        total_items = int(bag.get("totalItems") or bag.get("totalDishes") or len(dishes))
    except (TypeError, ValueError):
        total_items = len(dishes)

    ice = bool(bag.get("icePackRequired", False))
    addr2 = bag.get("address2", "") or ""
    city = bag.get("city", "") or ""
    prov = bag.get("province", "") or ""
    postal = bag.get("postalCode", "") or ""

    row = {
        "BAG_BARCODE": str(bag.get("bagBarcode", "")).strip(),
        "Shipping Name": bag.get("shippingName", "") or "",
        "Shipping Address 1": bag.get("address1", "") or "",
        "Shipping Address 2": addr2,
        "Shipping City": city,
        "Shipping Province": prov,
        "Shipping Postal Code": postal,
        "Shipping Phone": format_phone(bag.get("phone", "") or ""),
        "ZONE_NUMBER": zone,
        "HOUSEHOLD_MEMBERS": household,
        "DISH_LIST": dish_list,
        "TOTAL_ITEMS": total_items,
        "ICE_PACK_REQUIRED": ice,
    }
    row["line_shippingName"] = row["Shipping Name"]
    row["line_address"] = (
        f"{row['Shipping Address 1']}, {row['Shipping Address 2']}"
        if row["Shipping Address 2"]
        else row["Shipping Address 1"]
    )
    row["line_city"] = f"{city}, {prov} {postal}".strip(", ")
    row["line_zone"] = f"Zone {zone}" if zone else ""
    row["line_totalItems"] = f"Total {total_items} dish(es)"
    row["line_icePack"] = "Ice Pack Required" if ice else ""
    row["line_household"] = household
    row["line_dishes"] = dish_list
    return row


def generate_ppt_from_json(
    bags: list[dict],
    template_path: Path,
    output_path: Path,
) -> int:
    if not bags:
        raise ValueError("No bags to render")

    template_path = Path(template_path)
    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")

    template_path = Path(template_path)
    use_shipping_layout = not _uses_one_pager_template(template_path)

    prs = Presentation(str(template_path))
    template_slide = prs.slides[0]

    for bag in bags:
        row = json_bag_to_slide_row(bag)
        barcode = row["BAG_BARCODE"]
        if not barcode:
            print(f"⚠️ Skipping bag with no bagBarcode: {bag}")
            continue

        if use_shipping_layout:
            new_slide = copy_slide(template_slide, prs)
        else:
            new_slide = copy_slide_with_images(template_slide, prs)

        populate_bag_slide(new_slide, row, template_path)
        add_code128_barcode(new_slide, prs, barcode)

    r_id = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(r_id)
    del prs.slides._sldIdLst[0]

    save_path = Path(output_path)
    save_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(save_path))
    return len(prs.slides)


def print_summary(bags: list[dict], limit: int) -> None:
    dates = sorted({str(b.get("deliveryDate", "")) for b in bags if b.get("deliveryDate")})
    print(f"Bags in JSON: {len(bags)}")
    if dates:
        print(f"Delivery dates: {', '.join(dates)}")
    print(f"\nFirst {min(limit, len(bags))} bag(s):")
    for i, bag in enumerate(bags[:limit]):
        dishes = bag.get("dishes") or []
        print(
            f"  [{i}] {bag.get('bagBarcode')} — {bag.get('shippingName')} "
            f"(zone {bag.get('zone')}, {len(dishes)} dishes)"
        )


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Generate bag sticker PPT preview from bag_barcode_mapping.json",
    )
    parser.add_argument("--json", type=Path, help="Path to bag JSON (default: bag_barcode_mapping.json)")
    parser.add_argument(
        "-o",
        "--output",
        type=Path,
        help="Output .pptx path (default: preview_from_json_YYYYMMDD_HHMMSS.pptx)",
    )
    parser.add_argument("--template", type=Path, default=DEFAULT_BAG_TEMPLATE, help="PPT template path")
    parser.add_argument("--bag", help="Only render this bagBarcode")
    parser.add_argument("--limit", type=int, default=0, help="Max bags to render (0 = all)")
    parser.add_argument("--list-only", action="store_true", help="List JSON bags, do not build PPT")
    parser.add_argument("--find-json", action="store_true", help="List candidate JSON files and exit")
    parser.add_argument(
        "--mock-phone",
        action="store_true",
        help="Fill missing phone in JSON with 212-555-1XXX (preview only, does not save)",
    )
    args = parser.parse_args()

    if args.find_json:
        print("Bag mapping JSON files:")
        for p in find_json_candidates(BASE_DIR):
            print(f"  {p}")
        return 0

    json_path = args.json
    if json_path is None:
        candidates = find_json_candidates(BASE_DIR)
        json_path = DEFAULT_JSON if DEFAULT_JSON.exists() else (candidates[0] if candidates else None)
    if json_path is None:
        print(f"No JSON found. Put file at {DEFAULT_JSON} or pass --json PATH")
        return 1

    json_path = Path(json_path).resolve()
    print(f"Loading: {json_path}")
    mapping = load_mapping(json_path, inject_mock_phone=args.mock_phone)
    if args.mock_phone:
        print("Applied mock phone for entries missing phone (preview only).")

    if args.bag:
        mapping = [b for b in mapping if b.get("bagBarcode") == args.bag]
        if not mapping:
            print(f"Bag {args.bag!r} not found in JSON.")
            return 1

    if args.limit and args.limit > 0:
        mapping = mapping[: args.limit]

    print_summary(mapping, limit=10)

    if args.list_only:
        return 0

    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    output = args.output or (BASE_DIR / f"preview_from_json_{ts}.pptx")

    try:
        n = generate_ppt_from_json(mapping, args.template, output)
    except Exception as e:
        print(f"Failed to generate PPT: {e}")
        return 1

    print(f"\n✅ Wrote {n} slide(s) → {output}")
    print(f"   Template: {args.template}")
    print("   Open the file to review barcode size, address layout, and ZONE footer.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
